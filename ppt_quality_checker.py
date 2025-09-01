import os
import json
import time
import logging
from datetime import datetime
from typing import Dict, List, Optional, Tuple, Any
import hashlib
from pathlib import Path
import sys
import traceback
import re
from collections import Counter
import statistics
import streamlit as st

# Set page config as the very first Streamlit command
st.set_page_config(
    page_title="EduPPT Quality Analyzer Pro",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)


# Core imports with error handling
def check_dependencies():
    """Check and report missing dependencies"""
    missing_deps = []

    try:
        import pandas as pd
        import numpy as np
        import matplotlib.pyplot as plt
        import seaborn as sns
        from wordcloud import WordCloud
    except ImportError as e:
        missing_deps.append(f"Core dependencies: {e}")

    try:
        import nltk
        from nltk.data import find
        from nltk.tokenize import sent_tokenize, word_tokenize
        from nltk.corpus import stopwords
    except ImportError:
        missing_deps.append("NLTK not installed. Please run: pip install nltk")

    try:
        from langchain_community.document_loaders import UnstructuredPowerPointLoader
        from langchain_groq import ChatGroq
        from langchain_core.prompts import (
            SystemMessagePromptTemplate, HumanMessagePromptTemplate, ChatPromptTemplate
        )
        from langchain_core.output_parsers import StrOutputParser
    except ImportError as e:
        missing_deps.append(f"LangChain dependencies missing: {e}")

    try:
        from unstructured.partition.pptx import partition_pptx
        global UNSTRUCTURED_AVAILABLE
        UNSTRUCTURED_AVAILABLE = True
    except ImportError:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            UNSTRUCTURED_AVAILABLE = False
        except ImportError:
            missing_deps.append("No PowerPoint processing library available. Please install: pip install python-pptx")

    return missing_deps


# Store dependency check results
if 'dependency_check_done' not in st.session_state:
    st.session_state.dependency_errors = check_dependencies()
    st.session_state.dependency_check_done = True

# Import libraries after dependency check
try:
    import pandas as pd
    import numpy as np
    import matplotlib.pyplot as plt
    import seaborn as sns
    from wordcloud import WordCloud
    import nltk
    from nltk.data import find
    from nltk.tokenize import sent_tokenize, word_tokenize
    from nltk.corpus import stopwords
    from langchain_community.document_loaders import UnstructuredPowerPointLoader
    from langchain_groq import ChatGroq
    from langchain_core.prompts import (
        SystemMessagePromptTemplate, HumanMessagePromptTemplate, ChatPromptTemplate
    )
    from langchain_core.output_parsers import StrOutputParser

    try:
        from unstructured.partition.pptx import partition_pptx

        UNSTRUCTURED_AVAILABLE = True
    except ImportError:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        UNSTRUCTURED_AVAILABLE = False

except ImportError:
    pass

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger(__name__)


def safe_nltk_download():
    """Safely download NLTK data with error handling"""
    required_data = ['punkt', 'stopwords', 'averaged_perceptron_tagger', 'vader_lexicon']
    success = True
    messages = []

    for data_item in required_data:
        try:
            if data_item == 'punkt':
                find('tokenizers/punkt')
            elif data_item == 'stopwords':
                find('corpora/stopwords')
            elif data_item == 'averaged_perceptron_tagger':
                find('taggers/averaged_perceptron_tagger')
            elif data_item == 'vader_lexicon':
                find('vader_lexicon')
        except LookupError:
            try:
                nltk.download(data_item, quiet=True)
                messages.append(f"✅ NLTK {data_item} downloaded")
            except Exception as e:
                messages.append(f"⚠️ NLTK download warning for {data_item}: {e}")
                success = False

    return success, messages


class EduPPTQualityChecker:
    """Enhanced PowerPoint Quality Checker for Educational Content targeting B.Tech Students"""

    SUPPORTED_MODELS = {
        "llama-3.1-70b-versatile": "Llama 3.1 70B (Recommended - Best Quality)",
        "llama-3.1-8b-instant": "Llama 3.1 8B (Fast)",
        "mixtral-8x7b-32768": "Mixtral 8x7B (Good Balance)",
        "gemma2-9b-it": "Gemma 2 9B (Efficient)",
        "llama3-groq-70b-8192-tool-use-preview": "Llama 3 Groq 70B (Tool Use)"
    }

    # Enhanced quality categories for educational content
    QUALITY_CATEGORIES = {
        "educational_content": "Educational Content Quality",
        "student_engagement": "B.Tech Student Engagement",
        "technical_accuracy": "Technical Content Accuracy",
        "visual_learning": "Visual Learning Effectiveness",
        "knowledge_progression": "Learning Progression Structure",
        "accessibility": "Educational Accessibility",
        "professional_standards": "Industry Standards Compliance"
    }

    # Educational content keywords for B.Tech students
    EDUCATIONAL_KEYWORDS = {
        'engagement': ['example', 'practice', 'exercise', 'problem', 'solution', 'application', 'real-world',
                       'case study'],
        'technical': ['algorithm', 'method', 'approach', 'implementation', 'theory', 'principle', 'concept',
                      'framework'],
        'learning': ['definition', 'explain', 'understand', 'learn', 'remember', 'apply', 'analyze', 'evaluate'],
        'btech_specific': ['engineering', 'programming', 'code', 'design', 'system', 'software', 'hardware', 'network']
    }

    def __init__(self):
        """Initialize the educational quality checker"""
        try:
            self.session_id = self._generate_session_id()
            self.data_dir = Path("data")
            self.cache_dir = Path("cache")
            self.reports_dir = Path("quality_reports")
            self._ensure_directories()
            self._initialize_nltk()
            logger.info(f"EduPPTQualityChecker initialized with session ID: {self.session_id}")
        except Exception as e:
            logger.error(f"Initialization failed: {str(e)}")
            raise

    def _generate_session_id(self) -> str:
        """Generate unique session ID"""
        return hashlib.md5(f"{datetime.now().isoformat()}".encode()).hexdigest()[:8]

    def _ensure_directories(self):
        """Create necessary directories with error handling"""
        try:
            for directory in [self.data_dir, self.cache_dir, self.reports_dir]:
                directory.mkdir(parents=True, exist_ok=True)
                logger.info(f"Directory ensured: {directory}")
        except Exception as e:
            logger.error(f"Failed to create directories: {str(e)}")
            self.data_dir = Path(".")
            self.cache_dir = Path(".")
            self.reports_dir = Path(".")

    def _initialize_nltk(self):
        """Initialize NLTK with comprehensive error handling"""
        try:
            success, messages = safe_nltk_download()
            if not success:
                logger.warning("NLTK initialization had warnings but continuing...")
            if 'nltk_messages' not in st.session_state:
                st.session_state.nltk_messages = messages
        except Exception as e:
            logger.error(f"NLTK initialization failed: {str(e)}")
            if 'nltk_error' not in st.session_state:
                st.session_state.nltk_error = str(e)

    def save_uploaded_file(self, uploaded_file) -> Path:
        """Save uploaded file with comprehensive error handling"""
        try:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            clean_name = "".join(c for c in uploaded_file.name if c.isalnum() or c in (' ', '.', '_')).rstrip()
            filename = f"{timestamp}_{clean_name}"
            save_path = self.data_dir / filename

            with open(save_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            logger.info(f"File saved successfully: {save_path}")
            return save_path
        except Exception as e:
            logger.error(f"Failed to save file: {str(e)}")
            raise Exception(f"Could not save uploaded file: {str(e)}")

    def _extract_enhanced_content(self, file_path: Path) -> Dict[str, Any]:
        """Extract enhanced content analysis for educational presentations"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(str(file_path))
            enhanced_data = {
                'slides': [],
                'total_slides': len(prs.slides),
                'slide_layouts': [],
                'fonts_used': set(),
                'colors_used': set(),
                'images_count': 0,
                'charts_count': 0,
                'tables_count': 0,
                'shapes_count': 0,
                'educational_elements': {
                    'code_blocks': 0,
                    'formulas': 0,
                    'diagrams': 0,
                    'examples': 0,
                    'definitions': 0,
                    'questions': 0
                }
            }

            for i, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'slide_number': i,
                    'layout_name': slide.slide_layout.name if slide.slide_layout else "Unknown",
                    'text_content': [],
                    'word_count': 0,
                    'sentence_count': 0,
                    'bullet_points': 0,
                    'images': 0,
                    'charts': 0,
                    'tables': 0,
                    'shapes': 0,
                    'text_boxes': 0,
                    'educational_score': 0,
                    'technical_terms': 0,
                    'learning_indicators': 0
                }

                enhanced_data['slide_layouts'].append(slide.slide_layout.name if slide.slide_layout else "Unknown")

                for shape in slide.shapes:
                    enhanced_data['shapes_count'] += 1
                    slide_data['shapes'] += 1

                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        slide_data['text_content'].append(text_content)
                        slide_data['word_count'] += len(text_content.split())

                        # Enhanced educational content analysis
                        text_lower = text_content.lower()

                        # Count educational elements
                        if any(keyword in text_lower for keyword in ['example', 'for example', 'instance']):
                            enhanced_data['educational_elements']['examples'] += 1

                        if any(keyword in text_lower for keyword in ['definition', 'define', 'means', 'is defined as']):
                            enhanced_data['educational_elements']['definitions'] += 1

                        if '?' in text_content:
                            enhanced_data['educational_elements']['questions'] += text_content.count('?')

                        # Technical content detection
                        if any(keyword in text_lower for keyword in self.EDUCATIONAL_KEYWORDS['technical']):
                            slide_data['technical_terms'] += 1

                        # Learning progression indicators
                        if any(keyword in text_lower for keyword in self.EDUCATIONAL_KEYWORDS['learning']):
                            slide_data['learning_indicators'] += 1

                        # Code block detection (common in engineering presentations)
                        if re.search(r'[{}();=]', text_content) or 'algorithm' in text_lower:
                            enhanced_data['educational_elements']['code_blocks'] += 1

                        # Formula detection
                        if re.search(r'[∫∑∏αβγδλμπσΣΠΔΛΜ]|[a-zA-Z]\^[0-9]|[0-9]+\s*[+\-*/]\s*[0-9]', text_content):
                            enhanced_data['educational_elements']['formulas'] += 1

                        # Sentence analysis
                        try:
                            sentences = sent_tokenize(text_content)
                            slide_data['sentence_count'] += len(sentences)
                        except:
                            slide_data['sentence_count'] += text_content.count('.') + text_content.count(
                                '!') + text_content.count('?')

                        # Bullet points
                        slide_data['bullet_points'] += text_content.count('•') + text_content.count(
                            '*') + text_content.count('-')
                        slide_data['text_boxes'] += 1

                    # Shape type analysis
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        slide_data['images'] += 1
                        enhanced_data['images_count'] += 1
                        # Assume technical diagrams if in engineering context
                        enhanced_data['educational_elements']['diagrams'] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        slide_data['charts'] += 1
                        enhanced_data['charts_count'] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        slide_data['tables'] += 1
                        enhanced_data['tables_count'] += 1

                enhanced_data['slides'].append(slide_data)

            return enhanced_data

        except Exception as e:
            logger.error(f"Enhanced content extraction failed: {str(e)}")
            return {'slides': [], 'total_slides': 0, 'error': str(e)}

    def _analyze_educational_content(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze educational content quality specifically for B.Tech students"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            # Content structure analysis
            word_counts = [slide['word_count'] for slide in slides]
            avg_words = statistics.mean(word_counts) if word_counts else 0

            # Educational content metrics
            analysis['metrics']['average_words_per_slide'] = round(avg_words, 1)
            analysis['metrics']['total_educational_elements'] = sum(enhanced_data['educational_elements'].values())
            analysis['metrics']['examples_count'] = enhanced_data['educational_elements']['examples']
            analysis['metrics']['definitions_count'] = enhanced_data['educational_elements']['definitions']
            analysis['metrics']['questions_count'] = enhanced_data['educational_elements']['questions']
            analysis['metrics']['code_blocks'] = enhanced_data['educational_elements']['code_blocks']
            analysis['metrics']['formulas_count'] = enhanced_data['educational_elements']['formulas']

            # Learning progression analysis
            total_technical = sum(slide['technical_terms'] for slide in slides)
            total_learning_indicators = sum(slide['learning_indicators'] for slide in slides)

            analysis['metrics']['technical_density'] = round(total_technical / len(slides), 1)
            analysis['metrics']['learning_progression_score'] = round(total_learning_indicators / len(slides), 1)

            # Score calculation based on educational best practices
            score = 100

            # Optimal word count for educational content (slightly higher than business presentations)
            if avg_words > 60:
                score -= min(25, (avg_words - 60) * 0.4)
                analysis['issues'].append(f"High word count per slide ({avg_words:.1f} avg) - may overwhelm students")
                analysis['recommendations'].append(
                    "Reduce text density - aim for 30-50 words per slide for better comprehension")
            elif avg_words < 15:
                score -= 15
                analysis['issues'].append("Very low content density - may lack sufficient information")
                analysis['recommendations'].append("Add more explanatory content while maintaining clarity")

            # Educational element bonuses
            if enhanced_data['educational_elements']['examples'] > 0:
                score += min(15, enhanced_data['educational_elements']['examples'] * 3)
                analysis['recommendations'].append(
                    f"Excellent use of examples ({enhanced_data['educational_elements']['examples']})")
            else:
                score -= 10
                analysis['issues'].append("No practical examples found")
                analysis['recommendations'].append("Add real-world examples to improve understanding")

            if enhanced_data['educational_elements']['definitions'] > 0:
                score += min(10, enhanced_data['educational_elements']['definitions'] * 2)
            else:
                analysis['recommendations'].append("Consider adding clear definitions for technical terms")

            if enhanced_data['educational_elements']['questions'] < len(slides) * 0.2:
                score -= 15
                analysis['issues'].append("Insufficient interactive questions for student engagement")
                analysis['recommendations'].append("Add more questions to test understanding and maintain engagement")

            # Technical content appropriateness
            if total_technical / len(slides) < 0.5:
                score -= 10
                analysis['issues'].append("Low technical content density for B.Tech level")
                analysis['recommendations'].append("Increase technical depth appropriate for engineering students")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Educational content analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def _analyze_student_engagement(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze engagement factors specific to B.Tech students"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            # Engagement analysis
            total_questions = enhanced_data['educational_elements']['questions']
            interactive_slides = sum(
                1 for slide in slides if slide['charts'] > 0 or slide['images'] > 0 or slide['tables'] > 0)

            analysis['metrics']['questions_per_slide'] = round(total_questions / len(slides), 2)
            analysis['metrics']['interactive_slide_percentage'] = round((interactive_slides / len(slides)) * 100, 1)
            analysis['metrics']['visual_element_density'] = round(
                (enhanced_data['images_count'] + enhanced_data['charts_count']) / len(slides), 1)

            # Practical application indicators
            all_text = ' '.join([' '.join(slide['text_content']) for slide in slides])
            practical_keywords = ['practice', 'implement', 'build', 'create', 'develop', 'solve', 'hands-on', 'lab',
                                  'project']
            practical_score = sum(all_text.lower().count(keyword) for keyword in practical_keywords)
            analysis['metrics']['practical_application_score'] = practical_score

            # B.Tech specific engagement factors
            code_examples = enhanced_data['educational_elements']['code_blocks']
            analysis['metrics']['code_examples'] = code_examples

            # Score calculation
            score = 100

            # Question density (crucial for student engagement)
            if total_questions < len(slides) * 0.3:
                score -= 20
                analysis['issues'].append("Insufficient questions for active learning")
                analysis['recommendations'].append(
                    "Add more questions throughout the presentation (aim for 1 question per 3-4 slides)")

            # Visual engagement
            if interactive_slides / len(slides) < 0.4:
                score -= 15
                analysis['issues'].append("Low visual engagement - may lose student attention")
                analysis['recommendations'].append("Add more diagrams, charts, and visual aids")

            # Practical application
            if practical_score < 2:
                score -= 15
                analysis['issues'].append("Limited practical application references")
                analysis['recommendations'].append("Include more hands-on examples and practical applications")

            # Code examples for technical content
            if code_examples > 0:
                score += min(15, code_examples * 5)
                analysis['recommendations'].append(f"Great use of code examples ({code_examples})")
            else:
                analysis['recommendations'].append(
                    "Consider adding code snippets or pseudocode for better technical understanding")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Student engagement analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def _analyze_technical_accuracy(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze technical accuracy and depth for engineering students"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            all_text = ' '.join([' '.join(slide['text_content']) for slide in slides])

            # Technical terminology analysis
            technical_terms = 0
            for keyword_list in self.EDUCATIONAL_KEYWORDS.values():
                technical_terms += sum(all_text.lower().count(keyword) for keyword in keyword_list)

            analysis['metrics']['technical_term_frequency'] = round(technical_terms / len(all_text.split()) * 100,
                                                                    2) if all_text else 0
            analysis['metrics']['formula_count'] = enhanced_data['educational_elements']['formulas']
            analysis['metrics']['code_blocks'] = enhanced_data['educational_elements']['code_blocks']

            # Technical depth indicators
            complexity_indicators = ['implementation', 'architecture', 'optimization', 'complexity', 'performance']
            complexity_score = sum(all_text.lower().count(indicator) for indicator in complexity_indicators)
            analysis['metrics']['technical_complexity_score'] = complexity_score

            # Mathematical content (important for engineering)
            math_symbols = len(re.findall(r'[∫∑∏αβγδλμπσΣΠΔΛΜ]|[a-zA-Z]\^[0-9]', all_text))
            analysis['metrics']['mathematical_notation_count'] = math_symbols

            # Score calculation
            score = 100

            # Technical depth assessment
            if technical_terms < len(slides) * 2:
                score -= 20
                analysis['issues'].append("Insufficient technical depth for B.Tech level")
                analysis['recommendations'].append(
                    "Increase technical terminology and concepts appropriate for engineering students")

            # Formula and mathematical content
            if enhanced_data['educational_elements']['formulas'] == 0 and 'mathematics' in all_text.lower():
                score -= 10
                analysis['issues'].append("Mathematical content lacks proper formula representation")
                analysis['recommendations'].append("Use proper mathematical notation and formulas")

            # Code examples for programming topics
            if any(prog_term in all_text.lower() for prog_term in ['programming', 'algorithm', 'code', 'software']) and \
                    enhanced_data['educational_elements']['code_blocks'] == 0:
                score -= 15
                analysis['issues'].append("Programming content lacks code examples")
                analysis['recommendations'].append("Add code snippets, pseudocode, or algorithm representations")

            # Complexity appropriateness
            if complexity_score > 0:
                score += min(10, complexity_score * 2)
                analysis['recommendations'].append("Good coverage of technical complexity")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Technical accuracy analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def _analyze_visual_learning(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze visual learning effectiveness for engineering education"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            # Visual element analysis
            total_visual_elements = enhanced_data['images_count'] + enhanced_data['charts_count'] + enhanced_data[
                'tables_count']
            visual_density = total_visual_elements / len(slides)

            analysis['metrics']['visual_elements_per_slide'] = round(visual_density, 1)
            analysis['metrics']['diagram_count'] = enhanced_data['educational_elements']['diagrams']
            analysis['metrics']['chart_to_slide_ratio'] = round(enhanced_data['charts_count'] / len(slides), 2)

            # Text-visual balance
            text_only_slides = len(
                [slide for slide in slides if slide['images'] == 0 and slide['charts'] == 0 and slide['tables'] == 0])
            visual_balance_ratio = (len(slides) - text_only_slides) / len(slides)

            analysis['metrics']['text_only_slides'] = text_only_slides
            analysis['metrics']['visual_balance_percentage'] = round(visual_balance_ratio * 100, 1)

            # Learning style accommodation
            kinesthetic_elements = enhanced_data['educational_elements']['code_blocks'] + \
                                   enhanced_data['educational_elements']['examples']
            analysis['metrics']['hands_on_learning_elements'] = kinesthetic_elements

            # Score calculation
            score = 100

            # Visual density requirements for technical education
            if visual_density < 0.5:
                score -= 25
                analysis['issues'].append(f"Low visual content density ({visual_density:.1f} elements/slide)")
                analysis['recommendations'].append(
                    "Add more diagrams, flowcharts, and visual aids for better technical understanding")

            # Text-only slide penalty (higher for technical content)
            if text_only_slides / len(slides) > 0.5:
                score -= 30
                analysis['issues'].append(f"{text_only_slides} slides are text-only - not ideal for technical learning")
                analysis['recommendations'].append("Convert text-heavy slides to include visual representations")

            # Charts and graphs for data representation
            if enhanced_data['charts_count'] == 0 and len(slides) > 5:
                score -= 10
                analysis['issues'].append("No charts or graphs found - missing data visualization")
                analysis['recommendations'].append("Add charts, graphs, or data visualizations where applicable")

            # Reward for good visual learning design
            if 0.7 <= visual_density <= 1.5:
                score += 15
                analysis['recommendations'].append("Excellent visual learning balance")

            if kinesthetic_elements > len(slides) * 0.3:
                score += 10
                analysis['recommendations'].append("Good provision for hands-on learning elements")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Visual learning analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def _analyze_knowledge_progression(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze learning progression and knowledge building structure"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            # Progression analysis
            word_progression = [slide['word_count'] for slide in slides]
            complexity_progression = [slide['technical_terms'] + slide['learning_indicators'] for slide in slides]

            # Check for logical progression
            def check_progression_smoothness(values):
                if len(values) < 3:
                    return 1.0
                differences = [abs(values[i + 1] - values[i]) for i in range(len(values) - 1)]
                return 1.0 - (statistics.stdev(differences) / max(values) if max(values) > 0 else 0)

            content_smoothness = check_progression_smoothness(word_progression)
            complexity_smoothness = check_progression_smoothness(complexity_progression)

            analysis['metrics']['content_progression_smoothness'] = round(content_smoothness, 2)
            analysis['metrics']['complexity_progression_smoothness'] = round(complexity_smoothness, 2)

            # Learning structure indicators
            intro_slides = len([slide for slide in slides[:3] if any(keyword in ' '.join(slide['text_content']).lower()
                                                                     for keyword in
                                                                     ['introduction', 'overview', 'agenda',
                                                                      'objectives'])])
            conclusion_slides = len(
                [slide for slide in slides[-3:] if any(keyword in ' '.join(slide['text_content']).lower()
                                                       for keyword in ['conclusion', 'summary', 'recap', 'takeaway'])])

            analysis['metrics']['has_introduction_structure'] = intro_slides > 0
            analysis['metrics']['has_conclusion_structure'] = conclusion_slides > 0
            analysis['metrics']['structural_completeness'] = intro_slides + conclusion_slides

            # Score calculation
            score = 100

            # Progression smoothness
            if content_smoothness < 0.7:
                score -= 15
                analysis['issues'].append("Uneven content distribution affecting learning flow")
                analysis['recommendations'].append("Balance content distribution for smoother learning progression")

            if complexity_smoothness < 0.6:
                score -= 15
                analysis['issues'].append("Abrupt complexity changes may confuse students")
                analysis['recommendations'].append("Gradually increase complexity throughout the presentation")

            # Structure completeness
            if not intro_slides:
                score -= 10
                analysis['issues'].append("Missing clear introduction/overview")
                analysis['recommendations'].append("Add introduction slides with learning objectives")

            if not conclusion_slides:
                score -= 10
                analysis['issues'].append("Missing summary/conclusion")
                analysis['recommendations'].append("Add conclusion slides to reinforce key concepts")

            # Reward good progression
            if content_smoothness > 0.8 and complexity_smoothness > 0.7:
                score += 15
                analysis['recommendations'].append("Excellent learning progression structure")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Knowledge progression analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def _analyze_educational_accessibility(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze accessibility for diverse learning needs in engineering education"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            # Text complexity analysis
            all_text = ' '.join([' '.join(slide['text_content']) for slide in slides])
            sentences = sent_tokenize(all_text) if all_text else []

            if sentences:
                sentence_lengths = [len(word_tokenize(sentence)) for sentence in sentences]
                avg_sentence_length = statistics.mean(sentence_lengths)

                # Technical vocabulary complexity
                technical_ratio = sum(len(word) > 8 for word in word_tokenize(all_text.lower())) / len(
                    word_tokenize(all_text)) if all_text else 0
            else:
                avg_sentence_length = 0
                technical_ratio = 0

            analysis['metrics']['average_sentence_length'] = round(avg_sentence_length, 1)
            analysis['metrics']['technical_vocabulary_ratio'] = round(technical_ratio * 100, 1)

            # Multiple learning modalities
            visual_elements = enhanced_data['images_count'] + enhanced_data['charts_count']
            textual_elements = sum(slide['text_boxes'] for slide in slides)
            kinesthetic_elements = enhanced_data['educational_elements']['code_blocks'] + \
                                   enhanced_data['educational_elements']['examples']

            analysis['metrics']['visual_learning_support'] = visual_elements
            analysis['metrics']['textual_learning_support'] = textual_elements
            analysis['metrics']['kinesthetic_learning_support'] = kinesthetic_elements

            # Accessibility indicators
            heavy_text_slides = len([slide for slide in slides if slide['word_count'] > 80])
            analysis['metrics']['heavy_text_slides'] = heavy_text_slides

            # Score calculation
            score = 100

            # Sentence complexity for technical content
            if avg_sentence_length > 25:
                score -= 20
                analysis['issues'].append(
                    f"Long sentences ({avg_sentence_length:.1f} words avg) may hinder comprehension")
                analysis['recommendations'].append("Break down complex sentences for better accessibility")
            elif 12 <= avg_sentence_length <= 20:
                score += 10
                analysis['recommendations'].append("Good sentence length for technical content")

            # Technical vocabulary balance
            if technical_ratio > 0.4:
                score -= 15
                analysis['issues'].append("Very high technical vocabulary ratio may be overwhelming")
                analysis['recommendations'].append("Balance technical terms with explanations and simpler language")

            # Multiple learning modalities support
            modality_score = min(visual_elements, 1) + min(textual_elements, 1) + min(kinesthetic_elements, 1)
            if modality_score < 2:
                score -= 15
                analysis['issues'].append("Limited support for different learning styles")
                analysis['recommendations'].append(
                    "Include visual, textual, and hands-on elements to support diverse learners")
            else:
                score += modality_score * 5
                analysis['recommendations'].append("Good support for multiple learning modalities")

            # Heavy text slide penalty
            if heavy_text_slides > 0:
                score -= heavy_text_slides * 3
                analysis['issues'].append(f"{heavy_text_slides} slides have excessive text (>80 words)")
                analysis['recommendations'].append("Reduce text density for better cognitive load management")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Educational accessibility analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def _analyze_industry_standards(self, enhanced_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze compliance with educational industry standards"""
        try:
            analysis = {
                'overall_score': 0,
                'metrics': {},
                'issues': [],
                'recommendations': []
            }

            slides = enhanced_data.get('slides', [])
            if not slides:
                return analysis

            # Industry standard metrics for educational presentations
            slide_count = len(slides)
            analysis['metrics']['total_slides'] = slide_count

            # Content organization standards
            all_text = ' '.join([' '.join(slide['text_content']) for slide in slides])

            # Learning objectives identification
            objective_keywords = ['objective', 'goal', 'learn', 'understand', 'by the end', 'students will']
            objectives_present = any(keyword in all_text.lower() for keyword in objective_keywords)
            analysis['metrics']['learning_objectives_present'] = objectives_present

            # Assessment integration
            assessment_keywords = ['quiz', 'test', 'assignment', 'homework', 'practice', 'exercise']
            assessment_integration = sum(all_text.lower().count(keyword) for keyword in assessment_keywords)
            analysis['metrics']['assessment_integration_score'] = assessment_integration

            # Industry-relevant content
            industry_keywords = ['industry', 'professional', 'career', 'job', 'workplace', 'company', 'project']
            industry_relevance = sum(all_text.lower().count(keyword) for keyword in industry_keywords)
            analysis['metrics']['industry_relevance_score'] = industry_relevance

            # Timing considerations (educational presentations should be appropriately paced)
            estimated_time = slide_count * 2.5  # Assuming 2.5 minutes per slide for educational content
            analysis['metrics']['estimated_duration_minutes'] = round(estimated_time, 1)

            # Score calculation based on educational industry standards
            score = 100

            # Slide count optimization for educational content
            if slide_count < 8:
                score -= 15
                analysis['issues'].append("Presentation may be too brief for comprehensive learning")
                analysis['recommendations'].append("Expand content to cover topic more thoroughly")
            elif slide_count > 40:
                score -= 20
                analysis['issues'].append("Presentation may be too long for student attention span")
                analysis['recommendations'].append("Consider breaking into multiple sessions or modules")

            # Learning objectives requirement
            if not objectives_present:
                score -= 20
                analysis['issues'].append("No clear learning objectives identified")
                analysis['recommendations'].append("Add explicit learning objectives at the beginning")
            else:
                score += 10
                analysis['recommendations'].append("Clear learning objectives enhance educational effectiveness")

            # Assessment integration
            if assessment_integration == 0:
                score -= 10
                analysis['issues'].append("No assessment or practice elements found")
                analysis['recommendations'].append("Include practice questions, exercises, or assignments")
            elif assessment_integration > 0:
                score += min(15, assessment_integration * 3)
                analysis['recommendations'].append("Good integration of assessment elements")

            # Industry relevance for B.Tech students
            if industry_relevance == 0:
                score -= 10
                analysis['issues'].append("Limited connection to industry applications")
                analysis['recommendations'].append("Add industry examples and career relevance")
            elif industry_relevance > 0:
                score += min(10, industry_relevance * 2)
                analysis['recommendations'].append("Good industry relevance for career preparation")

            # Timing appropriateness
            if estimated_time > 90:
                score -= 15
                analysis['issues'].append(f"Estimated duration ({estimated_time:.1f} min) may exceed attention span")
                analysis['recommendations'].append("Consider shorter segments or break points")

            analysis['overall_score'] = max(0, min(100, score))
            return analysis

        except Exception as e:
            logger.error(f"Industry standards analysis failed: {str(e)}")
            return {'overall_score': 0, 'metrics': {}, 'issues': [str(e)], 'recommendations': []}

    def generate_enhanced_ai_feedback(self, enhanced_data: Dict[str, Any], quality_analysis: Dict[str, Any],
                                      api_key: str, model: str) -> str:
        """Generate AI-powered comprehensive feedback for educational presentations"""
        try:
            # Prepare enhanced context for educational AI analysis
            context = {
                'total_slides': enhanced_data.get('total_slides', 0),
                'educational_elements': enhanced_data.get('educational_elements', {}),
                'overall_scores': {category: analysis.get('overall_score', 0)
                                   for category, analysis in quality_analysis.items()},
                'key_issues': [],
                'key_recommendations': [],
                'target_audience': 'B.Tech Engineering Students'
            }

            for category, analysis in quality_analysis.items():
                context['key_issues'].extend(analysis.get('issues', []))
                context['key_recommendations'].extend(analysis.get('recommendations', []))

            # Create static prompt without template variables to avoid LangChain issues
            system_prompt = """You are an expert educational technology consultant specializing in creating 
            effective presentations for engineering students. You have extensive experience in B.Tech curriculum 
            design, student engagement strategies, and educational technology best practices. Your analysis focuses 
            on maximizing learning outcomes, student engagement, and knowledge retention for technical subjects."""

            # Build the user prompt with actual values (no template variables)
            educational_elements_str = str(context['educational_elements'])
            issues_text = '\n'.join('- ' + issue for issue in context['key_issues'][:12])
            recommendations_text = '\n'.join('- ' + rec for rec in context['key_recommendations'][:12])

            user_prompt = f"""Analyze this educational PowerPoint presentation designed for B.Tech engineering students:

PRESENTATION OVERVIEW:
- Total Slides: {context['total_slides']}
- Target Audience: {context['target_audience']}
- Educational Elements Found: {educational_elements_str}

QUALITY SCORES:
- Educational Content: {context['overall_scores'].get('educational_content', 0)}/100
- Student Engagement: {context['overall_scores'].get('student_engagement', 0)}/100  
- Technical Accuracy: {context['overall_scores'].get('technical_accuracy', 0)}/100
- Visual Learning: {context['overall_scores'].get('visual_learning', 0)}/100
- Knowledge Progression: {context['overall_scores'].get('knowledge_progression', 0)}/100
- Educational Accessibility: {context['overall_scores'].get('accessibility', 0)}/100
- Industry Standards: {context['overall_scores'].get('professional_standards', 0)}/100

IDENTIFIED ISSUES:
{issues_text}

CURRENT RECOMMENDATIONS:
{recommendations_text}

Please provide comprehensive feedback with:

1. Educational Effectiveness Assessment (3-4 sentences)
   - How well does this serve B.Tech students learning needs?
   - What is the overall pedagogical strength?

2. Top 3 Educational Strengths
   - Focus on what works well for engineering education
   - Highlight effective learning strategies used

3. Top 3 Critical Improvements for B.Tech Students
   - Prioritize changes that will most impact learning outcomes
   - Consider engineering student attention spans and learning preferences

4. Specific Action Plan (7-10 concrete steps)
   - Detailed, implementable improvements
   - Focus on enhancing student engagement and comprehension
   - Include industry-standard educational practices

5. Student Engagement Enhancement Strategies
   - Specific methods to increase participation
   - Ways to make technical content more accessible

6. Professional Learning Rating
   - Rate as: Excellent/Good/Developing/Needs Major Improvement
   - Justify the rating with specific educational criteria

7. Benchmarking Against Industry Standards
   - Compare to best practices in engineering education
   - Suggest alignment with current educational technology trends

Keep recommendations practical, specific, and focused on maximizing learning outcomes for engineering students."""

            # Initialize AI model with direct prompts
            llm = ChatGroq(
                model=model,
                api_key=api_key,
                temperature=0.2,
                max_tokens=3000,
                timeout=45
            )

            # Create messages directly without template variables
            from langchain_core.messages import SystemMessage, HumanMessage

            messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=user_prompt)
            ]

            # Generate response directly
            response = llm.invoke(messages)

            return response.content if hasattr(response, 'content') else str(response)

        except Exception as e:
            logger.error(f"Enhanced AI feedback generation failed: {str(e)}")
            return f"AI feedback unavailable: {str(e)}"

    def run_comprehensive_analysis(self, file_path: Path, use_ai: bool = False,
                                   api_key: str = "", model: str = "") -> Dict[str, Any]:
        """Run comprehensive educational quality analysis"""
        try:
            logger.info("Starting comprehensive educational PPT analysis")

            # Extract enhanced content
            enhanced_data = self._extract_enhanced_content(file_path)

            if 'error' in enhanced_data:
                raise Exception(enhanced_data['error'])

            # Run all educational quality analyses
            quality_analysis = {
                'educational_content': self._analyze_educational_content(enhanced_data),
                'student_engagement': self._analyze_student_engagement(enhanced_data),
                'technical_accuracy': self._analyze_technical_accuracy(enhanced_data),
                'visual_learning': self._analyze_visual_learning(enhanced_data),
                'knowledge_progression': self._analyze_knowledge_progression(enhanced_data),
                'accessibility': self._analyze_educational_accessibility(enhanced_data),
                'professional_standards': self._analyze_industry_standards(enhanced_data)
            }

            # Calculate weighted overall score (educational content and engagement weighted higher)
            weights = {
                'educational_content': 0.25,
                'student_engagement': 0.25,
                'technical_accuracy': 0.15,
                'visual_learning': 0.15,
                'knowledge_progression': 0.10,
                'accessibility': 0.05,
                'professional_standards': 0.05
            }

            weighted_score = sum(quality_analysis[category]['overall_score'] * weights[category]
                                 for category in weights.keys())

            # Generate enhanced AI feedback if requested
            ai_feedback = ""
            if use_ai and api_key:
                ai_feedback = self.generate_enhanced_ai_feedback(enhanced_data, quality_analysis, api_key, model)

            # Compile results
            results = {
                'overall_score': round(weighted_score, 1),
                'enhanced_data': enhanced_data,
                'quality_analysis': quality_analysis,
                'ai_feedback': ai_feedback,
                'analysis_timestamp': datetime.now().isoformat(),
                'file_name': file_path.name,
                'target_audience': 'B.Tech Engineering Students',
                'analysis_type': 'Educational Content Analysis'
            }

            logger.info(f"Educational analysis completed with overall score: {weighted_score}")
            return results

        except Exception as e:
            logger.error(f"Comprehensive educational analysis failed: {str(e)}")
            raise


def create_enhanced_educational_ui():
    """Create modern educational-focused UI"""
    try:
        st.markdown("""
        <style>
        .main-header {
            background: linear-gradient(135deg, #1E3A8A 0%, #7C3AED 50%, #EC4899 100%);
            padding: 2.5rem;
            border-radius: 20px;
            color: white;
            text-align: center;
            margin-bottom: 2rem;
            box-shadow: 0 10px 40px rgba(30, 58, 138, 0.3);
        }

        .edu-card {
            background: linear-gradient(145deg, #ffffff 0%, #f8fafc 100%);
            padding: 2rem;
            border-radius: 15px;
            box-shadow: 0 8px 30px rgba(0, 0, 0, 0.1);
            border-left: 5px solid #7C3AED;
            margin: 1.5rem 0;
            transition: all 0.3s ease;
        }

        .edu-card:hover {
            transform: translateY(-3px);
            box-shadow: 0 12px 35px rgba(124, 58, 237, 0.15);
        }

        .score-card-excellent {
            background: linear-gradient(135deg, #10B981 0%, #059669 100%);
            color: white;
            padding: 1.5rem;
            border-radius: 12px;
            text-align: center;
            margin: 1rem 0;
            box-shadow: 0 6px 20px rgba(16, 185, 129, 0.3);
        }

        .score-card-good {
            background: linear-gradient(135deg, #3B82F6 0%, #1D4ED8 100%);
            color: white;
            padding: 1.5rem;
            border-radius: 12px;
            text-align: center;
            margin: 1rem 0;
            box-shadow: 0 6px 20px rgba(59, 130, 246, 0.3);
        }

        .score-card-developing {
            background: linear-gradient(135deg, #F59E0B 0%, #D97706 100%);
            color: white;
            padding: 1.5rem;
            border-radius: 12px;
            text-align: center;
            margin: 1rem 0;
            box-shadow: 0 6px 20px rgba(245, 158, 11, 0.3);
        }

        .score-card-poor {
            background: linear-gradient(135deg, #EF4444 0%, #DC2626 100%);
            color: white;
            padding: 1.5rem;
            border-radius: 12px;
            text-align: center;
            margin: 1rem 0;
            box-shadow: 0 6px 20px rgba(239, 68, 68, 0.3);
        }

        .metric-highlight {
            background: linear-gradient(135deg, #F3E8FF 0%, #E0E7FF 100%);
            padding: 1rem;
            border-radius: 8px;
            border-left: 4px solid #7C3AED;
            margin: 0.5rem 0;
        }

        .edu-metric {
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding: 0.75rem;
            background: #F8FAFC;
            border-radius: 6px;
            margin: 0.5rem 0;
            border-left: 3px solid #7C3AED;
        }

        .improvement-badge {
            background: #FEF3C7;
            color: #92400E;
            padding: 0.25rem 0.75rem;
            border-radius: 20px;
            font-size: 0.8rem;
            font-weight: bold;
        }

        .strength-badge {
            background: #D1FAE5;
            color: #065F46;
            padding: 0.25rem 0.75rem;
            border-radius: 20px;
            font-size: 0.8rem;
            font-weight: bold;
        }

        .stButton > button {
            background: linear-gradient(135deg, #7C3AED 0%, #EC4899 100%);
            color: white;
            border: none;
            border-radius: 12px;
            padding: 1rem 2.5rem;
            font-weight: bold;
            font-size: 1.1rem;
            transition: all 0.3s ease;
            box-shadow: 0 4px 15px rgba(124, 58, 237, 0.3);
        }

        .stButton > button:hover {
            transform: translateY(-2px);
            box-shadow: 0 8px 25px rgba(124, 58, 237, 0.4);
        }
        </style>
        """, unsafe_allow_html=True)
    except Exception as e:
        logger.warning(f"Enhanced UI styling failed: {str(e)}")


def display_educational_quality_scores(results: Dict[str, Any]):
    """Display quality scores with educational focus"""
    try:
        st.markdown("### 📊 Educational Quality Assessment Dashboard")

        # Overall score with educational context
        overall_score = results['overall_score']

        if overall_score >= 85:
            score_class = "score-card-excellent"
            rating = "🎓 EXCELLENT - Industry Leading"
        elif overall_score >= 70:
            score_class = "score-card-good"
            rating = "📚 GOOD - Meets Standards"
        elif overall_score >= 50:
            score_class = "score-card-developing"
            rating = "⚡ DEVELOPING - Needs Enhancement"
        else:
            score_class = "score-card-poor"
            rating = "🔧 NEEDS MAJOR IMPROVEMENT"

        st.markdown(f"""
        <div class="{score_class}">
            <h2>Overall Educational Quality Score</h2>
            <h1>{overall_score}/100</h1>
            <p>{rating}</p>
            <p><strong>Target:</strong> B.Tech Engineering Students</p>
        </div>
        """, unsafe_allow_html=True)

        # Educational metrics overview
        enhanced_data = results['enhanced_data']
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric("📖 Total Slides", enhanced_data['total_slides'])
        with col2:
            examples = enhanced_data['educational_elements']['examples']
            st.metric("💡 Examples", examples, delta="Good" if examples > 2 else "Add More")
        with col3:
            questions = enhanced_data['educational_elements']['questions']
            st.metric("❓ Questions", questions, delta="Excellent" if questions > 5 else "Needs More")
        with col4:
            code_blocks = enhanced_data['educational_elements']['code_blocks']
            st.metric("💻 Code Elements", code_blocks, delta="Great" if code_blocks > 0 else "Missing")

        # Category scores with educational context
        st.markdown("### 📈 Detailed Quality Breakdown")

        col1, col2 = st.columns(2)
        categories = list(results['quality_analysis'].keys())

        for i, (category, analysis) in enumerate(results['quality_analysis'].items()):
            score = analysis['overall_score']
            col = col1 if i % 2 == 0 else col2

            if score >= 80:
                status_emoji = "🟢"
                badge_class = "strength-badge"
                status_text = "Excellent"
            elif score >= 65:
                status_emoji = "🟡"
                badge_class = "strength-badge"
                status_text = "Good"
            else:
                status_emoji = "🔴"
                badge_class = "improvement-badge"
                status_text = "Needs Work"

            with col:
                st.markdown(f"""
                <div class="edu-card">
                    <h4>{status_emoji} {EduPPTQualityChecker.QUALITY_CATEGORIES[category]}</h4>
                    <div style="display: flex; justify-content: space-between; align-items: center;">
                        <h2 style="margin: 0;">{score}/100</h2>
                        <span class="{badge_class}">{status_text}</span>
                    </div>
                </div>
                """, unsafe_allow_html=True)

    except Exception as e:
        st.error(f"Error displaying educational scores: {str(e)}")


def display_educational_insights(results: Dict[str, Any]):
    """Display educational-specific insights and recommendations"""
    try:
        st.markdown("### 🎯 Educational Insights & Recommendations")

        # Priority recommendations for educational improvement
        all_recommendations = []
        critical_issues = []

        for category, analysis in results['quality_analysis'].items():
            if analysis['overall_score'] < 60:  # Critical improvement areas
                critical_issues.extend(analysis.get('issues', []))
            all_recommendations.extend(analysis.get('recommendations', []))

        # Display critical issues first
        if critical_issues:
            st.markdown("#### 🚨 Critical Areas for Improvement")
            for i, issue in enumerate(critical_issues[:5], 1):
                st.markdown(f"""
                <div class="metric-highlight">
                    <strong>Priority {i}:</strong> {issue}
                </div>
                """, unsafe_allow_html=True)

        # Educational effectiveness summary
        enhanced_data = results['enhanced_data']
        edu_elements = enhanced_data['educational_elements']

        st.markdown("#### 📚 Educational Effectiveness Summary")

        col1, col2 = st.columns(2)

        with col1:
            st.markdown("**Learning Support Elements:**")
            for element, count in edu_elements.items():
                icon = {"examples": "💡", "definitions": "📖", "questions": "❓",
                        "code_blocks": "💻", "formulas": "🧮", "diagrams": "📊"}.get(element, "📌")
                status = "✅" if count > 0 else "❌"
                st.markdown(f"{status} {icon} {element.replace('_', ' ').title()}: {count}")

        with col2:
            st.markdown("**Engagement Factors:**")
            total_slides = enhanced_data['total_slides']

            # Calculate engagement ratios
            question_ratio = edu_elements['questions'] / total_slides if total_slides > 0 else 0
            example_ratio = edu_elements['examples'] / total_slides if total_slides > 0 else 0
            visual_ratio = (enhanced_data['images_count'] + enhanced_data[
                'charts_count']) / total_slides if total_slides > 0 else 0

            st.metric("Questions per Slide", f"{question_ratio:.2f}", delta="Target: 0.3+")
            st.metric("Examples per Slide", f"{example_ratio:.2f}", delta="Target: 0.4+")
            st.metric("Visual Elements per Slide", f"{visual_ratio:.2f}", delta="Target: 0.6+")

        # Detailed category analysis
        st.markdown("#### 🔍 Detailed Category Analysis")

        for category, analysis in results['quality_analysis'].items():
            with st.expander(
                    f"📋 {EduPPTQualityChecker.QUALITY_CATEGORIES[category]} - {analysis['overall_score']}/100"):

                # Create tabs for better organization
                tab1, tab2, tab3 = st.tabs(["Issues", "Recommendations", "Metrics"])

                with tab1:
                    if analysis['issues']:
                        st.markdown("**⚠️ Issues Identified:**")
                        for issue in analysis['issues']:
                            st.markdown(f"- 🔴 {issue}")
                    else:
                        st.success("✅ No significant issues found in this category!")

                with tab2:
                    if analysis['recommendations']:
                        st.markdown("**💡 Improvement Recommendations:**")
                        for rec in analysis['recommendations']:
                            st.markdown(f"- 💡 {rec}")

                with tab3:
                    if analysis['metrics']:
                        st.markdown("**📊 Detailed Metrics:**")
                        metrics_df = pd.DataFrame([
                            {"Metric": k.replace('_', ' ').title(), "Value": v}
                            for k, v in analysis['metrics'].items()
                        ])
                        st.dataframe(metrics_df, use_container_width=True)

    except Exception as e:
        st.error(f"Error displaying educational insights: {str(e)}")


def create_educational_visualizations(results: Dict[str, Any]):
    """Create enhanced visualizations for educational analysis"""
    try:
        st.markdown("### 📈 Educational Analysis Visualizations")

        # Enhanced visualization with educational focus
        categories = list(results['quality_analysis'].keys())
        scores = [results['quality_analysis'][cat]['overall_score'] for cat in categories]

        fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(16, 12))

        # 1. Educational Quality Radar Chart
        category_labels = [EduPPTQualityChecker.QUALITY_CATEGORIES[cat] for cat in categories]
        colors = ['#7C3AED', '#EC4899', '#10B981', '#F59E0B', '#EF4444', '#3B82F6', '#8B5CF6']

        bars = ax1.barh(range(len(categories)), scores, color=colors[:len(categories)])
        ax1.set_yticks(range(len(categories)))
        ax1.set_yticklabels([cat.replace('_', ' ').title() for cat in categories])
        ax1.set_xlabel('Quality Score')
        ax1.set_title('Educational Quality Scores by Category')
        ax1.set_xlim(0, 100)

        # Add score labels
        for bar, score in zip(bars, scores):
            ax1.text(bar.get_width() + 1, bar.get_y() + bar.get_height() / 2,
                     f'{score:.0f}', va='center', fontweight='bold')

        # 2. Learning Element Distribution
        edu_elements = results['enhanced_data']['educational_elements']
        element_names = list(edu_elements.keys())
        element_counts = list(edu_elements.values())

        wedges, texts, autotexts = ax2.pie(element_counts, labels=element_names, autopct='%1.1f%%',
                                           colors=colors[:len(element_names)], startangle=90)
        ax2.set_title('Educational Elements Distribution')

        # 3. Content Progression Analysis
        slides = results['enhanced_data']['slides']
        slide_numbers = [slide['slide_number'] for slide in slides]
        word_counts = [slide['word_count'] for slide in slides]
        technical_scores = [slide['technical_terms'] + slide['learning_indicators'] for slide in slides]

        ax3_twin = ax3.twinx()
        line1 = ax3.plot(slide_numbers, word_counts, 'o-', color='#7C3AED', linewidth=2, label='Word Count')
        line2 = ax3_twin.plot(slide_numbers, technical_scores, 's-', color='#EC4899', linewidth=2,
                              label='Technical Density')

        ax3.set_xlabel('Slide Number')
        ax3.set_ylabel('Word Count', color='#7C3AED')
        ax3_twin.set_ylabel('Technical Elements', color='#EC4899')
        ax3.set_title('Content and Complexity Progression')
        ax3.grid(True, alpha=0.3)

        # Combine legends
        lines1, labels1 = ax3.get_legend_handles_labels()
        lines2, labels2 = ax3_twin.get_legend_handles_labels()
        ax3.legend(lines1 + lines2, labels1 + labels2, loc='upper left')

        # 4. Educational Effectiveness Heatmap
        effectiveness_data = []
        for slide in slides:
            effectiveness_data.append([
                slide['word_count'],
                slide['technical_terms'],
                slide['learning_indicators'],
                slide['images'] + slide['charts'],
                slide['bullet_points']
            ])

        if effectiveness_data:
            effectiveness_df = pd.DataFrame(effectiveness_data,
                                            columns=['Words', 'Technical', 'Learning', 'Visual', 'Structure'])

            sns.heatmap(effectiveness_df.T, annot=True, cmap='viridis', ax=ax4,
                        xticklabels=[f'S{i + 1}' for i in range(len(slides))],
                        yticklabels=['Words', 'Technical', 'Learning', 'Visual', 'Structure'])
            ax4.set_title('Slide-by-Slide Educational Effectiveness')
            ax4.set_xlabel('Slide Number')

        plt.tight_layout()
        st.pyplot(fig)

        # Educational word cloud with focus on learning terms
        st.markdown("#### 🔤 Educational Content Word Cloud")
        all_text = ' '.join([' '.join(slide['text_content']) for slide in slides])
        if len(all_text) > 50:
            try:
                # Filter out common words and focus on educational/technical terms
                stop_words = set(stopwords.words('english'))
                educational_text = ' '.join([word for word in all_text.split()
                                             if word.lower() not in stop_words and len(word) > 3])

                wordcloud = WordCloud(width=1000, height=500, background_color='white',
                                      colormap='plasma', max_words=100).generate(educational_text)

                fig, ax = plt.subplots(figsize=(14, 7))
                ax.imshow(wordcloud, interpolation='bilinear')
                ax.axis('off')
                ax.set_title('Key Educational Terms and Concepts', fontsize=18, fontweight='bold')
                st.pyplot(fig)
            except Exception as e:
                st.warning(f"Could not generate educational word cloud: {str(e)}")

        # Learning progression visualization
        st.markdown("#### 📈 Learning Progression Analysis")

        if len(slides) > 1:
            progression_data = {
                'Slide': slide_numbers,
                'Content Density': word_counts,
                'Technical Depth': [slide['technical_terms'] for slide in slides],
                'Learning Indicators': [slide['learning_indicators'] for slide in slides],
                'Visual Support': [slide['images'] + slide['charts'] + slide['tables'] for slide in slides]
            }

            progression_df = pd.DataFrame(progression_data)
            st.line_chart(progression_df.set_index('Slide'))

    except Exception as e:
        st.error(f"Error creating educational visualizations: {str(e)}")


def generate_educational_report(results: Dict[str, Any], checker) -> str:
    """Generate comprehensive educational quality report"""
    try:
        timestamp = datetime.fromisoformat(results['analysis_timestamp']).strftime('%Y-%m-%d %H:%M:%S')

        report_content = f"""# Educational PowerPoint Quality Analysis Report

**📚 Course Material:** {results['file_name']}  
**🎯 Target Audience:** {results['target_audience']}  
**📅 Analysis Date:** {timestamp}  
**🏆 Overall Educational Quality Score:** {results['overall_score']}/100

## Executive Summary

**Educational Effectiveness Rating:** {'🌟 EXCELLENT' if results['overall_score'] >= 85 else '✅ GOOD' if results['overall_score'] >= 70 else '⚡ DEVELOPING' if results['overall_score'] >= 50 else '🔧 NEEDS MAJOR IMPROVEMENT'}

This presentation scored **{results['overall_score']}/100** across all educational quality dimensions, specifically evaluated for B.Tech engineering students.

## Educational Content Overview

| Metric | Value | Industry Standard |
|--------|--------|------------------|
| Total Slides | {results['enhanced_data']['total_slides']} | 15-25 for 45min session |
| Educational Examples | {results['enhanced_data']['educational_elements']['examples']} | 3+ recommended |
| Interactive Questions | {results['enhanced_data']['educational_elements']['questions']} | 5+ recommended |
| Code/Technical Elements | {results['enhanced_data']['educational_elements']['code_blocks']} | 2+ for technical topics |
| Visual Learning Aids | {results['enhanced_data']['images_count'] + results['enhanced_data']['charts_count']} | 60%+ slides should have visuals |

## Detailed Quality Assessment

| Educational Category | Score | Status | Priority |
|---------------------|-------|--------|----------|"""

        for category, analysis in results['quality_analysis'].items():
            score = analysis['overall_score']
            status = 'Excellent ✅' if score >= 80 else 'Good 🟡' if score >= 65 else 'Needs Improvement 🔴'
            priority = 'Low' if score >= 80 else 'Medium' if score >= 65 else 'HIGH'

            report_content += f"\n| {EduPPTQualityChecker.QUALITY_CATEGORIES[category]} | {score}/100 | {status} | {priority} |"

        # Add detailed category analysis
        report_content += "\n\n## Detailed Analysis by Category\n"

        for category, analysis in results['quality_analysis'].items():
            report_content += f"\n### {EduPPTQualityChecker.QUALITY_CATEGORIES[category]} ({analysis['overall_score']}/100)\n\n"

            if analysis['issues']:
                report_content += "**🚨 Critical Issues:**\n"
                for issue in analysis['issues']:
                    report_content += f"- ❌ {issue}\n"
                report_content += "\n"

            if analysis['recommendations']:
                report_content += "**💡 Improvement Recommendations:**\n"
                for rec in analysis['recommendations']:
                    report_content += f"- 🔧 {rec}\n"
                report_content += "\n"

        # Add AI feedback if available
        if results['ai_feedback']:
            report_content += f"\n## 🤖 AI Educational Expert Analysis\n\n{results['ai_feedback']}\n"

        # Add action plan
        report_content += "\n## 🎯 Priority Action Plan\n\n"

        # Collect high-priority recommendations
        priority_actions = []
        for category, analysis in results['quality_analysis'].items():
            if analysis['overall_score'] < 70:
                priority_actions.extend(analysis.get('recommendations', [])[:2])

        for i, action in enumerate(priority_actions[:8], 1):
            report_content += f"{i}. {action}\n"

        report_content += f"\n## 📊 Educational Metrics Summary\n\n"

        # Add comprehensive metrics
        enhanced_data = results['enhanced_data']
        total_words = sum(slide['word_count'] for slide in enhanced_data['slides'])
        avg_words = total_words / len(enhanced_data['slides']) if enhanced_data['slides'] else 0

        report_content += f"""
**Content Analysis:**
- Total Words: {total_words:,}
- Average Words per Slide: {avg_words:.1f}
- Educational Examples: {enhanced_data['educational_elements']['examples']}
- Definitions Provided: {enhanced_data['educational_elements']['definitions']}
- Interactive Questions: {enhanced_data['educational_elements']['questions']}
- Code Examples: {enhanced_data['educational_elements']['code_blocks']}
- Mathematical Formulas: {enhanced_data['educational_elements']['formulas']}

**Visual Learning Support:**
- Total Images/Diagrams: {enhanced_data['images_count']}
- Charts and Graphs: {enhanced_data['charts_count']}
- Tables: {enhanced_data['tables_count']}
- Visual Elements per Slide: {(enhanced_data['images_count'] + enhanced_data['charts_count'] + enhanced_data['tables_count']) / len(enhanced_data['slides']):.2f}

**Engagement Metrics:**
- Questions per Slide: {enhanced_data['educational_elements']['questions'] / len(enhanced_data['slides']):.2f}
- Examples per Slide: {enhanced_data['educational_elements']['examples'] / len(enhanced_data['slides']):.2f}
- Interactive Slide Percentage: {((len(enhanced_data['slides']) - len([s for s in enhanced_data['slides'] if s['images'] == 0 and s['charts'] == 0 and s['tables'] == 0])) / len(enhanced_data['slides']) * 100):.1f}%
"""

        report_content += f"\n---\n*Educational Quality Report generated by EduPPT Quality Analyzer Pro v3.0*  \n*Specialized for B.Tech Engineering Education*"

        return report_content

    except Exception as e:
        logger.error(f"Failed to generate educational report: {str(e)}")
        return f"Report generation failed: {str(e)}"


def main():
    """Enhanced main application function for educational presentations"""
    try:
        # Check for dependency errors first
        if st.session_state.dependency_errors:
            st.error("Missing Dependencies")
            for error in st.session_state.dependency_errors:
                st.error(error)
            st.info("Please install missing dependencies and restart the application.")
            st.stop()

        create_enhanced_educational_ui()

        # Enhanced header for educational focus
        st.markdown("""
        <div class="main-header">
            <h1>🎓 EduPPT Quality Analyzer Pro</h1>
            <p>Advanced presentation analysis for B.Tech engineering education</p>
            <p><strong>Maximize student engagement • Ensure technical accuracy • Meet industry standards</strong></p>
        </div>
        """, unsafe_allow_html=True)

        # Display NLTK setup messages
        if hasattr(st.session_state, 'nltk_messages') and st.session_state.nltk_messages:
            with st.expander("System Setup Status"):
                for msg in st.session_state.nltk_messages:
                    st.info(msg)

        # Initialize enhanced checker
        try:
            if 'edu_checker' not in st.session_state:
                with st.spinner("Initializing Educational Quality Analyzer..."):
                    st.session_state.edu_checker = EduPPTQualityChecker()
            checker = st.session_state.edu_checker
        except Exception as e:
            st.error(f"Application initialization failed: {str(e)}")
            st.stop()

        # Enhanced sidebar for educational settings
        with st.sidebar:
            st.header("Educational Analysis Configuration")

            # Target audience specification
            st.subheader("Target Audience")
            target_year = st.selectbox(
                "B.Tech Year Level:",
                ["1st Year (Foundation)", "2nd Year (Core)", "3rd Year (Advanced)", "4th Year (Specialization)",
                 "All Years"]
            )

            subject_area = st.selectbox(
                "Subject Area:",
                ["Computer Science", "Electronics", "Mechanical", "Civil", "Chemical", "General Engineering",
                 "Mathematics", "Mixed Topics"]
            )

            # Enhanced analysis categories
            st.subheader("Analysis Focus Areas")
            selected_categories = st.multiselect(
                "Select quality dimensions:",
                options=list(checker.QUALITY_CATEGORIES.keys()),
                default=list(checker.QUALITY_CATEGORIES.keys()),
                format_func=lambda x: checker.QUALITY_CATEGORIES[x]
            )

            # AI Analysis with educational specialization
            st.subheader("AI Educational Expert Analysis")
            enable_ai = st.checkbox("Enable AI Educational Expert Feedback", value=False)

            if enable_ai:
                groq_api_key = st.text_input(
                    "GROQ API Key",
                    type="password",
                    help="Required for AI educational analysis"
                )

                ai_model = st.selectbox(
                    "AI Model Selection",
                    options=list(checker.SUPPORTED_MODELS.keys()),
                    format_func=lambda x: checker.SUPPORTED_MODELS[x],
                    index=0  # Default to Llama 3.1 70B for better educational analysis
                )

                if groq_api_key:
                    st.success("AI Educational Expert Ready")
            else:
                groq_api_key = ""
                ai_model = "llama3-70b-8192"

            # Enhanced options
            st.subheader("Report Options")
            generate_report = st.checkbox("Generate Comprehensive Educational Report", value=True)
            create_visuals = st.checkbox("Create Educational Visualizations", value=True)
            include_benchmarks = st.checkbox("Include Industry Benchmarks", value=True)

        # Main content area
        col1, col2 = st.columns([2, 1])

        with col1:
            st.markdown("""
            <div class="edu-card">
                <h3>📁 Upload Your Educational Presentation</h3>
                <p>Upload your PowerPoint file (.pptx) for comprehensive educational quality analysis</p>
                <p><strong>Optimized for:</strong> B.Tech curriculum, technical content, student engagement</p>
            </div>
            """, unsafe_allow_html=True)

            uploaded_file = st.file_uploader(
                "Choose a PowerPoint file (.pptx)",
                type=["pptx"],
                help="Upload your educational presentation for analysis"
            )

            if uploaded_file is not None:
                # Enhanced file information
                file_size_mb = uploaded_file.size / (1024 * 1024)

                st.markdown("**📋 File Information:**")
                file_info_cols = st.columns(3)
                with file_info_cols[0]:
                    st.metric("File Name", uploaded_file.name)
                with file_info_cols[1]:
                    st.metric("Size", f"{file_size_mb:.2f} MB")
                with file_info_cols[2]:
                    st.metric("Type", "PowerPoint Presentation")

                # Enhanced analysis button
                if st.button("🚀 Start Educational Quality Analysis", type="primary"):
                    try:
                        # Save and process file
                        with st.spinner("Processing educational content..."):
                            save_path = checker.save_uploaded_file(uploaded_file)

                        # Enhanced progress tracking
                        progress_container = st.container()
                        with progress_container:
                            progress_bar = st.progress(0)
                            status_text = st.empty()

                            # Analysis steps with educational context
                            status_text.text("📖 Extracting educational content...")
                            progress_bar.progress(20)

                            status_text.text("🧠 Analyzing learning effectiveness...")
                            progress_bar.progress(40)

                            status_text.text("🎯 Evaluating student engagement...")
                            progress_bar.progress(60)

                            # Perform comprehensive analysis
                            results = checker.run_comprehensive_analysis(
                                save_path,
                                use_ai=enable_ai and bool(groq_api_key),
                                api_key=groq_api_key,
                                model=ai_model
                            )

                            progress_bar.progress(80)
                            status_text.text("📊 Generating educational insights...")

                            # Add context to results
                            results['target_year'] = target_year
                            results['subject_area'] = subject_area

                            progress_bar.progress(100)
                            status_text.text("Educational analysis complete!")

                            time.sleep(1)
                            status_text.empty()
                            progress_bar.empty()

                        # Store results
                        st.session_state.edu_analysis_results = results

                        st.success("Educational quality analysis completed successfully!")

                        # Display quick summary
                        overall_score = results['overall_score']
                        if overall_score >= 85:
                            st.balloons()
                            st.success(f"Outstanding educational quality! Score: {overall_score}/100")
                        elif overall_score >= 70:
                            st.success(
                                f"Good educational quality with room for improvement. Score: {overall_score}/100")
                        else:
                            st.warning(
                                f"Significant improvements needed for optimal learning outcomes. Score: {overall_score}/100")

                    except Exception as e:
                        st.error(f"Educational analysis failed: {str(e)}")
                        logger.error(f"Analysis error: {str(e)}")

        with col2:
            # Enhanced info panels for educational context
            st.markdown("""
            <div class="edu-card">
                <h3>🎯 Educational Analysis Features</h3>
                <ul>
                    <li><strong>Learning Effectiveness:</strong> Content structure, examples, definitions</li>
                    <li><strong>Student Engagement:</strong> Questions, interactivity, practical applications</li>
                    <li><strong>Technical Accuracy:</strong> Appropriate depth for B.Tech level</li>
                    <li><strong>Visual Learning:</strong> Diagrams, charts, multimedia support</li>
                    <li><strong>Knowledge Progression:</strong> Logical flow, complexity building</li>
                    <li><strong>Educational Accessibility:</strong> Multiple learning styles support</li>
                    <li><strong>Industry Standards:</strong> Best practices in engineering education</li>
                </ul>
            </div>
            """, unsafe_allow_html=True)

            st.markdown("""
            <div class="edu-card">
                <h3>🎓 B.Tech Student Success Tips</h3>
                <ul>
                    <li><strong>Content:</strong> 30-50 words per slide optimal</li>
                    <li><strong>Examples:</strong> Real-world applications essential</li>
                    <li><strong>Questions:</strong> 1 question per 3-4 slides</li>
                    <li><strong>Visuals:</strong> 60%+ slides should have graphics</li>
                    <li><strong>Code:</strong> Include pseudocode/examples</li>
                    <li><strong>Progression:</strong> Build complexity gradually</li>
                    <li><strong>Assessment:</strong> Include practice problems</li>
                </ul>
            </div>
            """, unsafe_allow_html=True)

            # Industry benchmarks info
            if include_benchmarks:
                st.markdown("""
                <div class="edu-card">
                    <h3>📈 Industry Benchmarks</h3>
                    <p><strong>Top Educational Institutions:</strong></p>
                    <ul>
                        <li>Examples: 4+ per presentation</li>
                        <li>Visual Elements: 70%+ slides</li>
                        <li>Questions: 8+ interactive elements</li>
                        <li>Duration: 45-60 minutes optimal</li>
                        <li>Technical Depth: Appropriate for year level</li>
                    </ul>
                </div>
                """, unsafe_allow_html=True)

        # Display results if available
        if 'edu_analysis_results' in st.session_state:
            results = st.session_state.edu_analysis_results

            st.markdown("---")

            # Enhanced quality scores display
            display_educational_quality_scores(results)

            # AI educational feedback
            if results.get('ai_feedback'):
                st.markdown("### 🤖 AI Educational Expert Analysis")
                st.markdown(f"""
                <div class="edu-card">
                    {results['ai_feedback']}
                </div>
                """, unsafe_allow_html=True)

            # Educational insights
            display_educational_insights(results)

            # Educational visualizations
            if create_visuals:
                create_educational_visualizations(results)

            # Generate and download enhanced report
            if generate_report:
                try:
                    with st.spinner("Generating comprehensive educational report..."):
                        report_content = generate_educational_report(results, checker)

                    # Enhanced download section
                    st.markdown("### 📥 Download Educational Quality Report")

                    col1, col2 = st.columns(2)
                    with col1:
                        st.download_button(
                            label="📄 Download Detailed Report (Markdown)",
                            data=report_content,
                            file_name=f"educational_quality_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
                            mime="text/markdown",
                            type="primary"
                        )

                    with col2:
                        # Generate JSON summary for further processing
                        json_summary = {
                            'overall_score': results['overall_score'],
                            'category_scores': {cat: analysis['overall_score'] for cat, analysis in
                                                results['quality_analysis'].items()},
                            'educational_elements': results['enhanced_data']['educational_elements'],
                            'target_audience': results['target_audience'],
                            'analysis_timestamp': results['analysis_timestamp']
                        }

                        st.download_button(
                            label="📊 Download Data (JSON)",
                            data=json.dumps(json_summary, indent=2),
                            file_name=f"edu_analysis_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                            mime="application/json"
                        )

                except Exception as e:
                    st.warning(f"Report generation failed: {str(e)}")

            # Educational improvement suggestions
            st.markdown("### 🎯 Quick Improvement Actions")

            improvement_cols = st.columns(3)

            # Get top recommendations from lowest scoring categories
            sorted_categories = sorted(results['quality_analysis'].items(),
                                       key=lambda x: x[1]['overall_score'])

            for i, (category, analysis) in enumerate(sorted_categories[:3]):
                with improvement_cols[i]:
                    st.markdown(f"""
                    <div class="metric-highlight">
                        <h4>Priority {i + 1}: {checker.QUALITY_CATEGORIES[category]}</h4>
                        <p><strong>Score:</strong> {analysis['overall_score']}/100</p>
                        <p><strong>Quick Fix:</strong> {analysis['recommendations'][0] if analysis['recommendations'] else 'Focus on this area'}</p>
                    </div>
                    """, unsafe_allow_html=True)

        # Enhanced footer with educational focus
        st.markdown("---")
        st.markdown("""
        <div style="text-align: center; color: #666; padding: 2rem; background: linear-gradient(135deg, #F8FAFC 0%, #E2E8F0 100%); border-radius: 15px; margin-top: 2rem;">
            <h3>🎓 EduPPT Quality Analyzer Pro v3.0</h3>
            <p><strong>Specialized for B.Tech Engineering Education</strong></p>
            <p>Helping educators create engaging, effective, and industry-standard presentations</p>
            <p><em>Built with ❤️ for educational excellence</em></p>
        </div>
        """, unsafe_allow_html=True)

    except Exception as e:
        st.error(f"Critical application error: {str(e)}")
        logger.error(f"Critical error in main: {str(e)}")


if __name__ == "__main__":
    main()
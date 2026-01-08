"""
================================================================================
EduPPT Quality Analyzer Ultimate - Complete Edition v7.8 (AWS PATH Fix)
================================================================================
The MOST COMPREHENSIVE PowerPoint Quality Analyzer for B.Tech Education

FEATURES:
✅ FIXED: "dirname/ls not found" Error (Explicit PATH Injection)
✅ UI FIXED: Perfect Dark/Light Theme Compatibility
✅ FIXED: LibreOffice "Exit Status 1" Error (Custom User Profile)
✅ .ENV File Support for API Keys
✅ Google Gemini 2.0 Flash Multimodal Analysis (Best Accuracy)
✅ AWS EC2 STABILITY PATCH (Memory & LibreOffice Fixes)
✅ AUTO-CLEANUP on Startup (Fix for File Accumulation)
✅ Groq Llama 4 Analysis (Alternative)
✅ 7 Quality Categories with Detailed Metrics
✅ Slide-by-Slide AI Recommendations

Author: Enhanced for PhysicsWallah EdTech
Version: 7.8.0 (PATH Env Fix + Theme UI + AWS Stability)
================================================================================
"""

import os
import json
import time
import logging
import base64
import io
import re
import tempfile
import subprocess
import random  # Required for Retry Logic
import shutil
import gc  # CRITICAL: For AWS Memory Management
from datetime import datetime
from typing import Dict, List, Optional, Tuple, Any, Union
import hashlib
from pathlib import Path
import sys
import traceback
from collections import Counter
import statistics
import asyncio
from concurrent.futures import ThreadPoolExecutor
import threading
from functools import lru_cache

# ============================================
# STREAMLIT CONFIG - MUST BE FIRST
# ============================================
import streamlit as st

st.set_page_config(
    page_title="EduPPT Quality Analyzer Ultimate",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)


# ============================================
# DEPENDENCY MANAGEMENT & ENV LOADING
# ============================================

def check_all_dependencies():
    """Comprehensive dependency check with detailed reporting"""
    missing_deps = []
    warnings = []
    install_commands = []

    # Environment Management
    try:
        from dotenv import load_dotenv
        load_dotenv()  # Load environment variables immediately
    except ImportError:
        missing_deps.append("python-dotenv not installed")
        install_commands.append("pip install python-dotenv")

    # Core data science
    try:
        import pandas as pd
        import numpy as np
    except ImportError as e:
        missing_deps.append(f"pandas/numpy: {e}")
        install_commands.append("pip install pandas numpy")

    # Visualization
    try:
        import matplotlib.pyplot as plt
        import seaborn as sns
    except ImportError as e:
        missing_deps.append(f"matplotlib/seaborn: {e}")
        install_commands.append("pip install matplotlib seaborn")

    try:
        from wordcloud import WordCloud
    except ImportError:
        missing_deps.append("wordcloud not installed")
        install_commands.append("pip install wordcloud")

    # NLP
    try:
        import nltk
    except ImportError:
        missing_deps.append("nltk not installed")
        install_commands.append("pip install nltk")

    # PowerPoint
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        missing_deps.append("python-pptx not installed")
        install_commands.append("pip install python-pptx")

    # Image processing
    try:
        from PIL import Image
    except ImportError:
        missing_deps.append("Pillow not installed")
        install_commands.append("pip install Pillow")

    # Google Gemini (Primary AI)
    try:
        import google.generativeai as genai
        GEMINI_AVAILABLE = True
    except ImportError:
        warnings.append("google-generativeai not installed (recommended for best accuracy)")
        install_commands.append("pip install google-generativeai")
        GEMINI_AVAILABLE = False

    # PDF to Image (Optional but recommended)
    try:
        from pdf2image import convert_from_path
        PDF2IMAGE_AVAILABLE = True
    except ImportError:
        warnings.append("pdf2image not installed (optional, for better slide rendering)")
        install_commands.append("pip install pdf2image")
        PDF2IMAGE_AVAILABLE = False

    # LangChain Groq (Alternative AI)
    try:
        from langchain_groq import ChatGroq
        from langchain_core.messages import SystemMessage, HumanMessage
        LANGCHAIN_AVAILABLE = True
    except ImportError:
        warnings.append("langchain-groq not installed (optional, for Groq models)")
        install_commands.append("pip install langchain-groq langchain-core")
        LANGCHAIN_AVAILABLE = False

    return {
        'missing': missing_deps,
        'warnings': warnings,
        'install_commands': install_commands,
        'gemini_available': GEMINI_AVAILABLE if 'GEMINI_AVAILABLE' in dir() else False,
        'pdf2image_available': PDF2IMAGE_AVAILABLE if 'PDF2IMAGE_AVAILABLE' in dir() else False,
        'langchain_available': LANGCHAIN_AVAILABLE if 'LANGCHAIN_AVAILABLE' in dir() else False
    }


# Run dependency check once
if 'dep_check' not in st.session_state:
    st.session_state.dep_check = check_all_dependencies()

# ============================================
# IMPORT ALL LIBRARIES
# ============================================

try:
    import pandas as pd
    import numpy as np
    import matplotlib.pyplot as plt
    import seaborn as sns
    from wordcloud import WordCloud
    import nltk
    from nltk.tokenize import sent_tokenize, word_tokenize
    from nltk.corpus import stopwords
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches, Pt
    from PIL import Image
    from dotenv import load_dotenv  # Ensure import
except ImportError as e:
    st.error(f"Critical import error: {e}")
    st.stop()

# Optional imports with flags
GEMINI_AVAILABLE = False
try:
    import google.generativeai as genai

    GEMINI_AVAILABLE = True
except ImportError:
    pass

PDF2IMAGE_AVAILABLE = False
try:
    from pdf2image import convert_from_path

    PDF2IMAGE_AVAILABLE = True
except ImportError:
    pass

LANGCHAIN_AVAILABLE = False
try:
    from langchain_groq import ChatGroq
    from langchain_core.messages import SystemMessage, HumanMessage

    LANGCHAIN_AVAILABLE = True
except ImportError:
    pass

# ============================================
# LOGGING CONFIGURATION
# ============================================

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger(__name__)


# ============================================
# NLTK INITIALIZATION
# ============================================

def safe_nltk_download():
    """Safely download required NLTK data"""
    required = ['punkt', 'stopwords', 'averaged_perceptron_tagger']
    messages = []

    for item in required:
        try:
            if 'punkt' in item:
                nltk.data.find(f'tokenizers/{item}')
            else:
                nltk.data.find(f'corpora/{item}')
        except LookupError:
            try:
                nltk.download(item, quiet=True)
                messages.append(f"✅ Downloaded {item}")
            except Exception as e:
                messages.append(f"⚠️ Could not download {item}")

    return messages


# Initialize NLTK once
if 'nltk_initialized' not in st.session_state:
    st.session_state.nltk_messages = safe_nltk_download()
    st.session_state.nltk_initialized = True


# ============================================
# MEMORY & AWS UTILS (NEW)
# ============================================

def clean_memory():
    """Aggressive memory cleanup for EC2"""
    gc.collect()


def resize_image_for_memory(image: Image.Image, max_dim: int = 1024) -> Image.Image:
    """Resize image to prevent OOM while keeping quality for AI"""
    try:
        if max(image.size) > max_dim:
            ratio = max_dim / max(image.size)
            new_size = tuple(int(dim * ratio) for dim in image.size)
            return image.resize(new_size, Image.Resampling.LANCZOS)
        return image
    except Exception:
        return image


# ============================================
# GEMINI MULTIMODAL ANALYZER (FIXED)
# ============================================

class GeminiMultimodalAnalyzer:
    """
    Advanced multimodal analyzer using Google Gemini
    Provides best-in-class image + text understanding
    Works for ANY content - universal analyzer
    """

    MODELS = {
        "gemini-2.0-flash": "Gemini 2.0 Flash (Best - Recommended)",
        "gemini-2.0-flash-lite": "Gemini 2.0 Flash Lite (Faster)",
        "gemini-1.5-pro": "Gemini 1.5 Pro (High Quality)",
        "gemini-1.5-flash": "Gemini 1.5 Flash (Fast)",
        "gemini-1.5-flash-8b": "Gemini 1.5 Flash 8B (Lightweight)",
    }

    def __init__(self, api_key: str, model: str = "gemini-2.0-flash"):
        """Initialize Gemini with API key and model"""
        self.api_key = api_key
        self.model_name = model

        genai.configure(api_key=api_key)

        self.generation_config = genai.GenerationConfig(
            temperature=0.2,  # Low for accuracy
            top_p=0.95,
            top_k=40,
            max_output_tokens=4096,
        )

        self.model = genai.GenerativeModel(
            model_name=model,
            generation_config=self.generation_config
        )

        logger.info(f"Gemini initialized: {model}")

    def _generate_with_retry(self, inputs, max_retries=5):
        """
        Internal helper to handle 429 Quota errors with Exponential Backoff
        """
        base_delay = 12  # Start with 12 seconds (Google typically asks for 10-15s)

        for attempt in range(max_retries):
            try:
                # Memory Check before generation
                clean_memory()
                return self.model.generate_content(inputs)
            except Exception as e:
                # Check for Quota/Resource Exhausted errors
                error_str = str(e)
                if "429" in error_str or "ResourceExhausted" in error_str:
                    if attempt == max_retries - 1:
                        # Failed all retries
                        logger.error(f"Failed after {max_retries} retries: {e}")
                        raise e

                        # Exponential Backoff: 12s, 24s, 48s... + Random Jitter
                    wait_time = (base_delay * (2 ** attempt)) + random.uniform(0, 2)

                    # Notify user via Streamlit and Console
                    msg = f"⚠️ Quota limit hit. Pausing for {int(wait_time)}s before retry {attempt + 1}/{max_retries}..."
                    print(msg)
                    st.toast(msg, icon="⏳")  # Show toast notification
                    time.sleep(wait_time)
                else:
                    # Raise other errors immediately
                    raise e
        return None

    def analyze_slide(self, image: Image.Image, slide_text: str,
                      slide_number: int, total_slides: int,
                      subject_area: str = "Engineering") -> Dict[str, Any]:
        """
        Comprehensive slide analysis - works for ANY content
        Analyzes both visual and textual elements together
        """
        try:
            # Optimize image size immediately for memory safety on EC2
            image = resize_image_for_memory(image)

            if image.mode != 'RGB':
                image = image.convert('RGB')

            # Universal analysis prompt - works for ANY content
            prompt = f"""You are an expert educational content reviewer. Analyze this slide image comprehensively.

SLIDE INFORMATION:
- Slide Number: {slide_number} of {total_slides}
- Subject Area: {subject_area}
- Target Audience: B.Tech Engineering Students
- Text Content: "{slide_text[:1500] if slide_text else 'No text extracted'}"

PERFORM A COMPLETE ANALYSIS:

1. **CONCEPTUAL ACCURACY** (Score 0-100):
   - Are all concepts, facts, and information correct?
   - Do diagrams/images accurately represent the concepts they illustrate?
   - Are there any factual errors or misconceptions?
   - Are labels, arrows, and connections logically correct?
   - Does the visual content match what the text describes?

2. **VISUAL-TEXT ALIGNMENT** (Score 0-100):
   - Do images support and match the textual content?
   - Are there any contradictions between visuals and text?
   - Are visual metaphors and analogies appropriate?
   - Do illustrations help explain the concept correctly?

3. **EDUCATIONAL EFFECTIVENESS** (Score 0-100):
   - Will students understand the concept from this slide?
   - Is the complexity appropriate for B.Tech level?
   - Is information presented clearly and logically?
   - Are there examples or applications shown?

4. **CONTENT QUALITY** (Score 0-100):
   - Is the content complete and sufficient?
   - Are technical terms used correctly?
   - Is there appropriate depth for the topic?
   - Are there any missing important points?

5. **DESIGN & LAYOUT** (Score 0-100):
   - Is the slide visually balanced?
   - Is text readable (size, contrast)?
   - Is there appropriate use of whitespace?
   - Are colors and formatting consistent?

6. **ENGAGEMENT POTENTIAL** (Score 0-100):
   - Does the slide capture attention?
   - Are there interactive elements or questions?
   - Will students stay engaged with this content?

RESPOND IN THIS EXACT FORMAT:

**SCORES:**
- Conceptual Accuracy: [0-100]
- Visual-Text Alignment: [0-100]
- Educational Effectiveness: [0-100]
- Content Quality: [0-100]
- Design & Layout: [0-100]
- Engagement Potential: [0-100]
- Overall Score: [0-100]

**CRITICAL_ERRORS:**
[List ALL factual errors, conceptual mistakes, or misleading information. Be VERY specific about what is wrong and why. If none, write "None found."]

**CONCEPTUAL_ISSUES:**
[Describe any issues where visuals don't correctly represent concepts, or where there are logical problems]

**VISUAL_ISSUES:**
[Describe any design, layout, image quality, or visual clarity problems]

**TEXT_ISSUES:**
[Describe any text problems - accuracy, clarity, completeness, grammar, spelling]

**ALIGNMENT_PROBLEMS:**
[Describe any mismatches between what images show and what text says]

**IMMEDIATE_FIXES:**
[List the TOP 3 most important fixes needed, in order of priority]

**DETAILED_RECOMMENDATIONS:**
[Provide specific, actionable recommendations to improve this slide]

**STUDENT_LEARNING_IMPACT:**
[How will the current state of this slide affect student understanding?]

**POSITIVE_ASPECTS:**
[What is done well on this slide?]

**PRIORITY_LEVEL:** [CRITICAL/HIGH/MEDIUM/LOW]
(CRITICAL = Major errors that will mislead students)
(HIGH = Significant issues affecting learning)
(MEDIUM = Improvements needed but content is usable)
(LOW = Minor refinements suggested)

Be thorough, specific, and objective. This analysis should work for ANY educational content."""

            # Send to Gemini with RETRY logic
            response = self._generate_with_retry([prompt, image])

            # Parse response
            result = self._parse_response(response.text, slide_number)
            result['raw_response'] = response.text
            result['model'] = self.model_name

            return result

        except Exception as e:
            logger.error(f"Gemini analysis error for slide {slide_number}: {e}")
            return self._fallback_result(slide_number, str(e))

    def _parse_response(self, text: str, slide_number: int) -> Dict[str, Any]:
        """Parse Gemini response into structured data"""
        result = {
            'slide_number': slide_number,
            'conceptual_accuracy_score': 50,
            'visual_text_alignment_score': 50,
            'educational_effectiveness_score': 50,
            'content_quality_score': 50,
            'design_layout_score': 50,
            'engagement_score': 50,
            'overall_score': 50,
            'critical_errors': [],
            'conceptual_issues': '',
            'visual_issues': '',
            'text_issues': '',
            'alignment_problems': '',
            'immediate_fixes': [],
            'detailed_recommendations': '',
            'student_impact': '',
            'positive_aspects': '',
            'priority_level': 'MEDIUM'
        }

        try:
            lines = text.strip().split('\n')
            current_section = None
            section_content = []

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Extract scores
                score_patterns = [
                    ('Conceptual Accuracy', 'conceptual_accuracy_score'),
                    ('Visual-Text Alignment', 'visual_text_alignment_score'),
                    ('Educational Effectiveness', 'educational_effectiveness_score'),
                    ('Content Quality', 'content_quality_score'),
                    ('Design & Layout', 'design_layout_score'),
                    ('Engagement Potential', 'engagement_score'),
                    ('Overall Score', 'overall_score'),
                ]

                for pattern, key in score_patterns:
                    if pattern.lower() in line.lower():
                        score = re.search(r'(\d+)', line)
                        if score:
                            result[key] = min(100, max(0, int(score.group(1))))

                # Extract sections
                section_map = {
                    '**CRITICAL_ERRORS:**': 'critical_errors',
                    '**CONCEPTUAL_ISSUES:**': 'conceptual_issues',
                    '**VISUAL_ISSUES:**': 'visual_issues',
                    '**TEXT_ISSUES:**': 'text_issues',
                    '**ALIGNMENT_PROBLEMS:**': 'alignment_problems',
                    '**IMMEDIATE_FIXES:**': 'immediate_fixes',
                    '**DETAILED_RECOMMENDATIONS:**': 'detailed_recommendations',
                    '**STUDENT_LEARNING_IMPACT:**': 'student_impact',
                    '**POSITIVE_ASPECTS:**': 'positive_aspects',
                    '**PRIORITY_LEVEL:**': 'priority_level',
                }

                section_found = False
                for header, key in section_map.items():
                    if header in line:
                        # Save previous section
                        if current_section and section_content:
                            if current_section in ['critical_errors', 'immediate_fixes']:
                                result[current_section] = [s for s in section_content if s and len(s) > 3]
                            else:
                                result[current_section] = '\n'.join(section_content)

                        current_section = key
                        section_content = []
                        section_found = True

                        # Check for inline content
                        remaining = line.replace(header, '').strip()
                        if remaining:
                            section_content.append(remaining)
                        break

                if not section_found and current_section:
                    clean = re.sub(r'^\*+|\*+$', '', line).strip()
                    clean = clean.lstrip('-•').strip()
                    if clean and clean.lower() not in ['none', 'none found', 'none found.', 'n/a']:
                        section_content.append(clean)

            # Save last section
            if current_section and section_content:
                if current_section in ['critical_errors', 'immediate_fixes']:
                    result[current_section] = [s for s in section_content if s and len(s) > 3]
                elif current_section == 'priority_level':
                    content = ' '.join(section_content).upper()
                    if 'CRITICAL' in content:
                        result['priority_level'] = 'CRITICAL'
                    elif 'HIGH' in content:
                        result['priority_level'] = 'HIGH'
                    elif 'LOW' in content:
                        result['priority_level'] = 'LOW'
                    else:
                        result['priority_level'] = 'MEDIUM'
                else:
                    result[current_section] = '\n'.join(section_content)

            # Calculate overall if not set
            if result['overall_score'] == 50:
                scores = [
                    result['conceptual_accuracy_score'],
                    result['visual_text_alignment_score'],
                    result['educational_effectiveness_score'],
                    result['content_quality_score'],
                    result['design_layout_score'],
                    result['engagement_score']
                ]
                result['overall_score'] = int(statistics.mean(scores))

        except Exception as e:
            logger.error(f"Parse error: {e}")

        return result

    def _fallback_result(self, slide_number: int, error: str) -> Dict[str, Any]:
        """Generate fallback when analysis fails"""
        return {
            'slide_number': slide_number,
            'conceptual_accuracy_score': 50,
            'visual_text_alignment_score': 50,
            'educational_effectiveness_score': 50,
            'content_quality_score': 50,
            'design_layout_score': 50,
            'engagement_score': 50,
            'overall_score': 50,
            'critical_errors': [f'Analysis failed: {error}'],
            'conceptual_issues': 'Unable to analyze',
            'visual_issues': 'Unable to analyze',
            'text_issues': 'Unable to analyze',
            'alignment_problems': 'Unable to analyze',
            'immediate_fixes': ['Retry with valid API key'],
            'detailed_recommendations': 'Please retry analysis',
            'student_impact': 'Unknown',
            'positive_aspects': 'Unable to determine',
            'priority_level': 'MEDIUM',
            'error': error
        }


# ============================================
# GROQ/LLAMA ANALYZER (Alternative)
# ============================================

class GroqLlamaAnalyzer:
    """
    Alternative analyzer using Groq's Llama models
    Good for text-based analysis when Gemini is not available
    """

    MODELS = {
        "meta-llama/llama-4-scout-17b-16e-instruct": "Llama 4 Scout 17B (Multimodal)",
        "meta-llama/llama-4-maverick-17b-128e-instruct": "Llama 4 Maverick 17B (Multimodal)",
        "llama-3.3-70b-versatile": "Llama 3.3 70B (Text Only - High Quality)",
        "llama-3.1-8b-instant": "Llama 3.1 8B (Fast)",
        "mixtral-8x7b-32768": "Mixtral 8x7B (Balanced)",
    }

    VISION_MODELS = {
        "meta-llama/llama-4-scout-17b-16e-instruct",
        "meta-llama/llama-4-maverick-17b-128e-instruct"
    }

    def __init__(self, api_key: str, model: str = "llama-3.3-70b-versatile"):
        self.api_key = api_key
        self.model_name = model
        self.is_vision = model in self.VISION_MODELS

        self.llm = ChatGroq(
            model=model,
            api_key=api_key,
            temperature=0.3,
            max_tokens=2000,
            timeout=60,
            max_retries=2
        )

        logger.info(f"Groq initialized: {model}, Vision: {self.is_vision}")

    def analyze_slide(self, slide_text: str, slide_data: Dict,
                      slide_number: int, total_slides: int,
                      image_base64: str = None) -> Dict[str, Any]:
        """Analyze slide with Groq/Llama"""
        try:
            # Clean memory before calling LLM
            clean_memory()

            prompt = f"""Analyze this educational slide for B.Tech students:

SLIDE {slide_number}/{total_slides}
TEXT CONTENT: "{slide_text[:2000]}"

SLIDE METRICS:
- Word Count: {slide_data.get('word_count', 0)}
- Images: {slide_data.get('images', 0)}
- Charts: {slide_data.get('charts', 0)}
- Tables: {slide_data.get('tables', 0)}

Provide analysis in this format:

**OVERALL_SCORE:** [0-100]
**PRIORITY_LEVEL:** [CRITICAL/HIGH/MEDIUM/LOW]

**CRITICAL_ERRORS:**
[List any factual or conceptual errors]

**IMMEDIATE_FIXES:**
[Top 3 fixes needed]

**RECOMMENDATIONS:**
[Detailed improvement suggestions]

**POSITIVE_ASPECTS:**
[What is done well]"""

            if self.is_vision and image_base64:
                content = [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"}
                    }
                ]
                messages = [HumanMessage(content=content)]
            else:
                messages = [HumanMessage(content=prompt)]

            response = self.llm.invoke(messages)
            text = response.content if hasattr(response, 'content') else str(response)

            return self._parse_response(text, slide_number)

        except Exception as e:
            logger.error(f"Groq analysis error: {e}")
            return {
                'slide_number': slide_number,
                'overall_score': 50,
                'priority_level': 'MEDIUM',
                'critical_errors': [str(e)],
                'immediate_fixes': [],
                'recommendations': '',
                'positive_aspects': '',
                'error': str(e)
            }

    def _parse_response(self, text: str, slide_number: int) -> Dict[str, Any]:
        """Parse Groq response"""
        result = {
            'slide_number': slide_number,
            'overall_score': 50,
            'priority_level': 'MEDIUM',
            'critical_errors': [],
            'immediate_fixes': [],
            'recommendations': '',
            'positive_aspects': ''
        }

        # Extract score
        score_match = re.search(r'\*\*OVERALL_SCORE:\*\*\s*(\d+)', text)
        if score_match:
            result['overall_score'] = min(100, int(score_match.group(1)))

        # Extract priority
        if 'CRITICAL' in text.upper():
            result['priority_level'] = 'CRITICAL'
        elif 'HIGH' in text.upper():
            result['priority_level'] = 'HIGH'
        elif 'LOW' in text.upper():
            result['priority_level'] = 'LOW'

        # Extract sections (simplified)
        result['raw_response'] = text

        return result


# ============================================
# SLIDE IMAGE CONVERTER (AWS PATH PATCHED)
# ============================================

class SlideImageConverter:
    """Convert PowerPoint slides to images for analysis - AWS/EC2 Compatible"""

    def __init__(self, temp_dir: Path = None):
        self.temp_dir = temp_dir or Path(tempfile.mkdtemp())
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def convert_pptx_to_images(self, pptx_path: Path) -> List[Tuple[int, Image.Image, str]]:
        """
        Convert PPTX to images
        Returns: [(slide_number, image, text), ...]
        """
        results = []

        # Method 1: LibreOffice (best quality) - WITH EC2 PATH DETECTION
        try:
            results = self._convert_via_libreoffice(pptx_path)
            if results:
                return results
        except Exception as e:
            logger.warning(f"LibreOffice failed: {e}")
            st.toast("⚠️ LibreOffice failed or not found. Using fallback extractor...", icon="⚠️")

        # Method 2: Extract images from shapes (Fallback 1)
        try:
            logger.info("Falling back to shape extraction")
            results = self._convert_via_shapes(pptx_path)
            if results:
                return results
        except Exception as e:
            logger.warning(f"Shape extraction failed: {e}")

        # Method 3: Create text-based placeholder images (Last Resort)
        logger.info("Falling back to placeholders")
        results = self._create_placeholder_images(pptx_path)
        return results

    def _convert_via_libreoffice(self, pptx_path: Path) -> List[Tuple[int, Image.Image, str]]:
        """Convert using LibreOffice -> PDF -> Images (Cross-Platform: Mac + AWS/Linux)"""
        results = []
        pdf_path = self.temp_dir / f"{pptx_path.stem}.pdf"

        # 1. Find LibreOffice - CROSS-PLATFORM paths (Linux/AWS + Mac)
        soffice_candidates = [
            # Linux/AWS paths (Ubuntu, Debian, EC2)
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",

            # Mac paths
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/usr/local/bin/soffice",
            "/opt/homebrew/bin/soffice",

            # Fallback: system PATH search
            shutil.which("libreoffice"),
            shutil.which("soffice")
        ]

        soffice_path = None
        for candidate in soffice_candidates:
            if candidate and os.path.exists(str(candidate)) and os.access(str(candidate), os.X_OK):
                soffice_path = candidate
                logger.info(f"✅ Found LibreOffice at: {soffice_path}")
                break

        if not soffice_path:
            raise Exception(
                "LibreOffice not found. Install with: sudo apt-get install libreoffice (Linux) or brew install libreoffice (Mac)")

        # 2. Detect platform and use appropriate command
        import platform
        is_linux = platform.system() == "Linux"

        if is_linux:
            # AWS/Linux: Use custom user profile (required for headless server)
            user_profile_dir = Path(tempfile.mkdtemp(prefix="libreoffice_"))

            cmd = [
                str(soffice_path),
                f"--env:UserInstallation=file://{user_profile_dir.as_posix()}",
                "--headless",
                "--invisible",
                "--nocrashreport",
                "--nodefault",
                "--nofirststartwizard",
                "--nolockcheck",
                "--nologo",
                "--norestore",
                "--convert-to", "pdf",
                "--outdir", str(self.temp_dir.absolute()),
                str(pptx_path.absolute())
            ]
        else:
            # Mac: Simplified command (custom profile causes issues on Mac)
            user_profile_dir = None

            cmd = [
                str(soffice_path),
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.temp_dir.absolute()),
                str(pptx_path.absolute())
            ]

        try:
            logger.info(f"🔄 Running: {' '.join(cmd)}")

            # CRITICAL FIX: Manually inject PATH to environment
            # This fixes "dirname: not found" and "ls: not found" errors on EC2
            env = os.environ.copy()

            # Ensure standard paths are present
            current_path = env.get('PATH', '')
            required_paths = ['/usr/bin', '/bin', '/usr/local/bin']

            for p in required_paths:
                if p not in current_path:
                    current_path = f"{p}:{current_path}"

            env['PATH'] = current_path

            if is_linux and user_profile_dir:
                env['HOME'] = str(user_profile_dir)

            # Run conversion
            result = subprocess.run(
                cmd,
                capture_output=True,
                timeout=180,
                env=env,  # Pass the corrected environment
                cwd=str(Path.home()) if not is_linux else str(self.temp_dir)
            )

            # Log output
            stdout = result.stdout.decode('utf-8', errors='ignore')
            stderr = result.stderr.decode('utf-8', errors='ignore')

            logger.info(f"LibreOffice stdout: {stdout}")
            if stderr:
                logger.info(f"LibreOffice stderr: {stderr}")

            if result.returncode != 0:
                error_msg = f"Exit code {result.returncode}"
                if stderr:
                    error_msg += f": {stderr}"
                if stdout:
                    error_msg += f" | stdout: {stdout}"

                # Platform-specific error hints
                if is_linux and "Couldn't create" in stderr:
                    error_msg += " | TIP: Ensure LibreOffice is fully installed with: sudo apt-get install -y libreoffice libreoffice-core fonts-liberation"
                elif not is_linux and "Couldn't create pipe" in stderr:
                    error_msg += " | TIP: Run 'sudo xattr -r -d com.apple.quarantine /Applications/LibreOffice.app'"

                raise Exception(f"LibreOffice conversion failed: {error_msg}")

        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out (180s)")

        except Exception as e:
            raise Exception(f"LibreOffice execution error: {str(e)}")

        finally:
            # Cleanup temp profile (Linux only)
            if is_linux and user_profile_dir:
                try:
                    shutil.rmtree(user_profile_dir, ignore_errors=True)
                except:
                    pass

        # 4. Verify PDF was created
        if not pdf_path.exists():
            # Check for PDF with any name
            possible_pdfs = list(self.temp_dir.glob("*.pdf"))
            if possible_pdfs:
                pdf_path = possible_pdfs[0]
                logger.info(f"📄 Found PDF: {pdf_path}")
            else:
                files = list(self.temp_dir.glob("*"))
                logger.error(f"Files in temp_dir: {files}")
                raise Exception(f"PDF not created. Expected: {pdf_path}")

        # 5. Convert PDF to images
        if not PDF2IMAGE_AVAILABLE:
            raise Exception("pdf2image not installed. Run: pip install pdf2image")

        try:
            logger.info(f"🖼️ Converting PDF to images...")
            images = convert_from_path(
                str(pdf_path),
                dpi=150,
                thread_count=1 if is_linux else 2,  # Single thread on AWS
                fmt='jpeg'
            )
            logger.info(f"✅ Extracted {len(images)} images")

        except Exception as e:
            raise Exception(f"PDF to image conversion failed: {str(e)}")

        # 6. Extract text and combine
        prs = Presentation(str(pptx_path))

        for i, (img, slide) in enumerate(zip(images, prs.slides), 1):
            img = resize_image_for_memory(img)
            text = self._extract_text(slide)
            results.append((i, img, text))

        # 7. Cleanup PDF
        try:
            os.remove(pdf_path)
        except:
            pass

        return results

    def _convert_via_shapes(self, pptx_path: Path) -> List[Tuple[int, Image.Image, str]]:
        """Extract and composite shapes into slide images (Pure Python Fallback)"""
        results = []
        prs = Presentation(str(pptx_path))

        # Use smaller canvas for fallback to save memory
        slide_width = 800
        slide_height = 600

        for slide_num, slide in enumerate(prs.slides, 1):
            # Create base image
            img = Image.new('RGB', (slide_width, slide_height), '#1E3A5F')
            text = self._extract_text(slide)

            # Composite images from slide
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        if hasattr(shape, "image"):
                            shape_img = Image.open(io.BytesIO(shape.image.blob))

                            # Resize large shape images immediately
                            shape_img.thumbnail((400, 400))

                            # Handle transparency
                            if shape_img.mode == 'RGBA':
                                bg = Image.new('RGB', shape_img.size, '#1E3A5F')
                                bg.paste(shape_img, mask=shape_img.split()[3])
                                shape_img = bg
                            elif shape_img.mode != 'RGB':
                                shape_img = shape_img.convert('RGB')

                            # Center paste for rough approximation
                            x_offset = random.randint(50, slide_width - 450)
                            y_offset = random.randint(50, slide_height - 350)
                            img.paste(shape_img, (x_offset, y_offset))

                    except Exception as e:
                        logger.warning(f"Shape error: {e}")

            results.append((slide_num, img, text))

        return results

    def _create_placeholder_images(self, pptx_path: Path) -> List[Tuple[int, Image.Image, str]]:
        """Create placeholder images when conversion fails"""
        results = []
        prs = Presentation(str(pptx_path))

        for slide_num, slide in enumerate(prs.slides, 1):
            # Blue placeholder
            img = Image.new('RGB', (800, 600), '#1E3A5F')
            text = self._extract_text(slide)
            results.append((slide_num, img, text))

        return results

    def _extract_text(self, slide) -> str:
        """Extract all text from slide"""
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        return ' '.join(texts)


# ============================================
# COMPREHENSIVE QUALITY ANALYZER
# ============================================

class UltimateEduPPTAnalyzer:
    """
    COMPLETE PowerPoint Quality Analyzer
    Includes ALL features from previous versions PLUS enhanced AI analysis
    Works for ANY educational content
    """

    # Quality categories with descriptions
    QUALITY_CATEGORIES = {
        "conceptual_accuracy": "Conceptual & Factual Accuracy",
        "visual_alignment": "Visual-Text Alignment",
        "educational_content": "Educational Content Quality",
        "student_engagement": "Student Engagement",
        "technical_depth": "Technical Depth",
        "visual_design": "Visual Design & Layout",
        "learning_progression": "Learning Progression"
    }

    # Educational keywords for analysis
    EDUCATIONAL_KEYWORDS = {
        'engagement': ['example', 'practice', 'exercise', 'problem', 'solution',
                       'application', 'real-world', 'case study', 'demonstration',
                       'hands-on', 'interactive', 'activity', 'try', 'experiment'],
        'technical': ['algorithm', 'method', 'approach', 'implementation', 'theory',
                      'principle', 'concept', 'framework', 'architecture', 'design',
                      'pattern', 'optimization', 'complexity', 'efficiency'],
        'learning': ['definition', 'explain', 'understand', 'learn', 'remember',
                     'apply', 'analyze', 'evaluate', 'create', 'synthesize',
                     'compare', 'contrast', 'describe', 'identify'],
        'btech': ['engineering', 'programming', 'code', 'software', 'hardware',
                  'network', 'database', 'protocol', 'interface', 'module',
                  'component', 'system', 'design', 'development'],
        'assessment': ['quiz', 'test', 'assignment', 'homework', 'project',
                       'lab', 'experiment', 'report', 'presentation'],
        'industry': ['industry', 'professional', 'career', 'job', 'workplace',
                     'company', 'enterprise', 'business', 'real-world']
    }

    def __init__(self):
        """Initialize the complete analyzer"""
        self.session_id = hashlib.md5(
            f"{datetime.now().isoformat()}".encode()
        ).hexdigest()[:8]

        self.data_dir = Path("data")
        self.cache_dir = Path("cache")
        self.reports_dir = Path("quality_reports")
        self.images_dir = Path("slide_images")

        self._ensure_directories()

        # Initialize converters
        self.image_converter = SlideImageConverter(self.images_dir)

        # Cache for performance
        self._analysis_cache = {}

        logger.info(f"Analyzer initialized: {self.session_id}")

    def _ensure_directories(self):
        """Create required directories"""
        for d in [self.data_dir, self.cache_dir, self.reports_dir, self.images_dir]:
            d.mkdir(parents=True, exist_ok=True)

    def save_uploaded_file(self, uploaded_file) -> Path:
        """Save uploaded file and return path"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        clean_name = "".join(
            c for c in uploaded_file.name
            if c.isalnum() or c in (' ', '.', '_')
        ).rstrip()

        filename = f"{timestamp}_{clean_name}"
        save_path = self.data_dir / filename

        with open(save_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        logger.info(f"Saved: {save_path}")
        return save_path

    def extract_presentation_data(self, file_path: Path) -> Dict[str, Any]:
        """Extract comprehensive data from presentation"""
        try:
            prs = Presentation(str(file_path))

            data = {
                'total_slides': len(prs.slides),
                'slides': [],
                'educational_elements': {
                    'examples': 0,
                    'definitions': 0,
                    'questions': 0,
                    'code_blocks': 0,
                    'formulas': 0,
                    'diagrams': 0
                },
                'images_count': 0,
                'charts_count': 0,
                'tables_count': 0,
                'shapes_count': 0,
                'fonts_used': set(),
                'slide_layouts': []
            }

            for i, slide in enumerate(prs.slides, 1):
                slide_data = self._analyze_single_slide(slide, i)
                data['slides'].append(slide_data)

                # Aggregate
                data['images_count'] += slide_data['images']
                data['charts_count'] += slide_data['charts']
                data['tables_count'] += slide_data['tables']
                data['shapes_count'] += slide_data['shapes']
                data['slide_layouts'].append(slide_data['layout_name'])

                # Count educational elements
                text_lower = ' '.join(slide_data['text_content']).lower()

                if any(kw in text_lower for kw in ['example', 'for example', 'instance', 'case study']):
                    data['educational_elements']['examples'] += 1

                if any(kw in text_lower for kw in ['definition', 'define', 'means', 'is defined as']):
                    data['educational_elements']['definitions'] += 1

                data['educational_elements']['questions'] += text_lower.count('?')

                # Code detection
                if re.search(r'[{}();=]|def |class |function|algorithm', text_lower):
                    data['educational_elements']['code_blocks'] += 1

                # Formula detection
                if re.search(r'[∫∑∏αβγδλμπσ]|[a-zA-Z]\^[0-9]|\d+\s*[+\-*/]\s*\d+',
                             ' '.join(slide_data['text_content'])):
                    data['educational_elements']['formulas'] += 1

                if slide_data['images'] > 0 or slide_data['charts'] > 0:
                    data['educational_elements']['diagrams'] += 1

            # Convert set to list for JSON serialization
            data['fonts_used'] = list(data['fonts_used'])

            return data

        except Exception as e:
            logger.error(f"Extraction error: {e}")
            return {'error': str(e), 'slides': [], 'total_slides': 0}

    def _analyze_single_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """Detailed analysis of a single slide"""
        slide_data = {
            'slide_number': slide_num,
            'layout_name': slide.slide_layout.name if slide.slide_layout else "Unknown",
            'text_content': [],
            'word_count': 0,
            'sentence_count': 0,
            'bullet_points': 0,
            'images': 0,
            'charts': 0,
            'tables': 0,
            'shapes': 0,
            'text_boxes': 0,
            'technical_terms': 0,
            'learning_indicators': 0,
            'engagement_elements': 0,
            'content_density': 'medium',
            'slide_purpose': 'content',
            'complexity_level': 'medium',
            'basic_score': 70
        }

        for shape in slide.shapes:
            slide_data['shapes'] += 1

            # Text analysis
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_data['text_content'].append(text)

                try:
                    words = word_tokenize(text)
                    sentences = sent_tokenize(text)
                    slide_data['word_count'] += len(words)
                    slide_data['sentence_count'] += len(sentences)
                except:
                    slide_data['word_count'] += len(text.split())
                    slide_data['sentence_count'] += text.count('.') + text.count('!') + text.count('?')

                slide_data['bullet_points'] += text.count('•') + text.count('-') + text.count('*')
                slide_data['text_boxes'] += 1

                # Educational indicators
                text_lower = text.lower()
                for category, keywords in self.EDUCATIONAL_KEYWORDS.items():
                    for kw in keywords:
                        if kw in text_lower:
                            if category == 'technical':
                                slide_data['technical_terms'] += 1
                            elif category == 'learning':
                                slide_data['learning_indicators'] += 1
                            elif category == 'engagement':
                                slide_data['engagement_elements'] += 1

            # Shape type analysis
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_data['images'] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                slide_data['charts'] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                slide_data['tables'] += 1

        # Classify slide
        slide_data['content_density'] = self._classify_density(slide_data['word_count'])
        slide_data['slide_purpose'] = self._classify_purpose(
            ' '.join(slide_data['text_content']), slide_num,
            st.session_state.get('total_slides', 20)
        )
        slide_data['complexity_level'] = self._classify_complexity(
            slide_data['technical_terms'], slide_data['word_count']
        )
        slide_data['basic_score'] = self._calculate_basic_score(slide_data)

        return slide_data

    def _classify_density(self, word_count: int) -> str:
        if word_count < 20:
            return 'low'
        elif word_count < 60:
            return 'medium'
        else:
            return 'high'

    def _classify_purpose(self, text: str, slide_num: int, total: int) -> str:
        text_lower = text.lower()

        if slide_num <= 2 or any(kw in text_lower for kw in
                                 ['introduction', 'overview', 'agenda', 'objectives', 'outline']):
            return 'title'

        if slide_num >= total - 1 or any(kw in text_lower for kw in
                                         ['conclusion', 'summary', 'recap', 'takeaway', 'thank you']):
            return 'summary'

        if len(text.split()) < 15 and any(kw in text_lower for kw in
                                          ['next', 'moving on', 'let us', 'now we']):
            return 'transition'

        return 'content'

    def _classify_complexity(self, tech_terms: int, word_count: int) -> str:
        if word_count == 0:
            return 'low'
        ratio = tech_terms / word_count
        if ratio < 0.1:
            return 'low'
        elif ratio < 0.3:
            return 'medium'
        else:
            return 'high'

    def _calculate_basic_score(self, slide_data: Dict) -> int:
        """Calculate basic slide score before AI analysis"""
        score = 70

        # Word count optimization (30-50 is ideal)
        wc = slide_data['word_count']
        if 30 <= wc <= 50:
            score += 15
        elif 20 <= wc <= 60:
            score += 10
        elif wc > 100:
            score -= 25
        elif wc < 10 and slide_data['images'] == 0:
            score -= 20

        # Visual elements
        visuals = slide_data['images'] + slide_data['charts'] + slide_data['tables']
        if visuals > 0:
            score += min(15, visuals * 5)
        elif slide_data['word_count'] > 40:
            score -= 10  # Text-heavy without visuals

        # Engagement
        if slide_data['engagement_elements'] > 0:
            score += min(10, slide_data['engagement_elements'] * 3)

        # Technical depth
        if slide_data['technical_terms'] > 0:
            score += min(10, slide_data['technical_terms'] * 2)

        # Learning indicators
        if slide_data['learning_indicators'] > 0:
            score += min(10, slide_data['learning_indicators'] * 3)

        return max(0, min(100, score))

    def run_gemini_analysis(self, file_path: Path, presentation_data: Dict,
                            api_key: str, model: str) -> Dict[int, Dict[str, Any]]:
        """Run comprehensive Gemini multimodal analysis"""
        results = {}

        try:
            # Initialize Gemini
            analyzer = GeminiMultimodalAnalyzer(api_key, model)

            # Convert slides to images
            st.info("🔄 Converting slides to images...")
            slide_images = self.image_converter.convert_pptx_to_images(file_path)

            if not slide_images:
                st.error("Failed to convert slides")
                return {}

            st.success(f"✅ Converted {len(slide_images)} slides")

            # Analyze each slide
            progress = st.progress(0)
            status = st.empty()

            subject = st.session_state.get('subject_area', 'Engineering')
            total = len(slide_images)

            for i, (slide_num, image, text) in enumerate(slide_images):
                # Clean memory before every slide analysis
                clean_memory()

                status.text(f"🔍 Analyzing slide {slide_num}/{total} with Gemini...")

                # Get additional text from presentation data
                if slide_num <= len(presentation_data.get('slides', [])):
                    slide_info = presentation_data['slides'][slide_num - 1]
                    full_text = ' '.join(slide_info.get('text_content', [])) or text
                else:
                    full_text = text

                try:
                    analysis = analyzer.analyze_slide(
                        image, full_text, slide_num, total, subject
                    )

                    # Merge with basic data
                    if slide_num <= len(presentation_data.get('slides', [])):
                        analysis['slide_data'] = presentation_data['slides'][slide_num - 1]

                    results[slide_num] = analysis

                except Exception as e:
                    logger.error(f"Slide {slide_num} error: {e}")
                    results[slide_num] = analyzer._fallback_result(slide_num, str(e))

                progress.progress((i + 1) / total)

                # Explicitly delete image reference to free memory
                del image

            progress.empty()
            status.empty()

            st.success(f"✅ Gemini analysis complete: {len(results)} slides")

        except Exception as e:
            logger.error(f"Gemini analysis failed: {e}")
            st.error(f"Analysis error: {str(e)}")

        return results

    def run_groq_analysis(self, presentation_data: Dict, api_key: str,
                          model: str) -> Dict[int, Dict[str, Any]]:
        """Run Groq/Llama analysis as alternative"""
        results = {}

        try:
            analyzer = GroqLlamaAnalyzer(api_key, model)

            slides = presentation_data.get('slides', [])
            total = len(slides)

            progress = st.progress(0)
            status = st.empty()

            for i, slide_data in enumerate(slides):
                clean_memory()
                slide_num = slide_data['slide_number']
                status.text(f"🔍 Analyzing slide {slide_num}/{total} with Groq...")

                text = ' '.join(slide_data.get('text_content', []))

                try:
                    analysis = analyzer.analyze_slide(
                        text, slide_data, slide_num, total
                    )
                    analysis['slide_data'] = slide_data
                    results[slide_num] = analysis

                except Exception as e:
                    logger.error(f"Slide {slide_num} error: {e}")
                    results[slide_num] = {
                        'slide_number': slide_num,
                        'overall_score': slide_data.get('basic_score', 50),
                        'error': str(e)
                    }

                progress.progress((i + 1) / total)
                time.sleep(0.3)

            progress.empty()
            status.empty()

            st.success(f"✅ Groq analysis complete: {len(results)} slides")

        except Exception as e:
            logger.error(f"Groq analysis failed: {e}")
            st.error(f"Analysis error: {str(e)}")

        return results

    def calculate_quality_scores(self, presentation_data: Dict,
                                 ai_results: Dict) -> Dict[str, Dict[str, Any]]:
        """Calculate comprehensive quality scores across all categories"""
        scores = {}
        slides = presentation_data.get('slides', [])

        # 1. CONCEPTUAL ACCURACY (from AI)
        if ai_results:
            concept_scores = [
                r.get('conceptual_accuracy_score', r.get('overall_score', 50))
                for r in ai_results.values()
            ]
            critical_errors = []
            for sn, r in ai_results.items():
                for err in r.get('critical_errors', []):
                    if err and len(err) > 5:
                        critical_errors.append(f"Slide {sn}: {err}")

            scores['conceptual_accuracy'] = {
                'overall_score': int(statistics.mean(concept_scores)) if concept_scores else 50,
                'issues': critical_errors[:10],
                'recommendations': ['Review slides with critical errors'] if critical_errors else [],
                'metrics': {
                    'slides_with_errors': len([r for r in ai_results.values() if r.get('critical_errors')]),
                    'total_analyzed': len(ai_results)
                }
            }
        else:
            scores['conceptual_accuracy'] = {
                'overall_score': 50,
                'issues': ['AI analysis not performed'],
                'recommendations': ['Enable Gemini or Groq analysis'],
                'metrics': {}
            }

        # 2. VISUAL-TEXT ALIGNMENT (from AI)
        if ai_results:
            align_scores = [
                r.get('visual_text_alignment_score', r.get('overall_score', 50))
                for r in ai_results.values()
            ]
            align_issues = []
            for sn, r in ai_results.items():
                if r.get('alignment_problems') and len(r['alignment_problems']) > 5:
                    align_issues.append(f"Slide {sn}: {r['alignment_problems'][:200]}")

            scores['visual_alignment'] = {
                'overall_score': int(statistics.mean(align_scores)) if align_scores else 50,
                'issues': align_issues[:10],
                'recommendations': [],
                'metrics': {}
            }
        else:
            scores['visual_alignment'] = {
                'overall_score': 50,
                'issues': [],
                'recommendations': [],
                'metrics': {}
            }

        # 3. EDUCATIONAL CONTENT
        edu = presentation_data.get('educational_elements', {})
        edu_score = 100
        edu_issues = []
        edu_recs = []

        if edu.get('examples', 0) == 0:
            edu_score -= 20
            edu_issues.append("No practical examples found")
            edu_recs.append("Add real-world examples")
        elif edu.get('examples', 0) >= 3:
            edu_score += 5

        if edu.get('questions', 0) == 0:
            edu_score -= 15
            edu_issues.append("No interactive questions")
            edu_recs.append("Add questions to test understanding")

        if edu.get('definitions', 0) == 0:
            edu_score -= 10
            edu_issues.append("No clear definitions")
            edu_recs.append("Define key technical terms")

        if edu.get('code_blocks', 0) == 0 and 'programming' in st.session_state.get('subject_area', '').lower():
            edu_score -= 15
            edu_issues.append("No code examples for programming topic")
            edu_recs.append("Add code snippets and examples")

        scores['educational_content'] = {
            'overall_score': max(0, min(100, edu_score)),
            'issues': edu_issues,
            'recommendations': edu_recs,
            'metrics': edu
        }

        # 4. STUDENT ENGAGEMENT
        if slides:
            total_engagement = sum(s.get('engagement_elements', 0) for s in slides)
            visual_slides = sum(1 for s in slides if s.get('images', 0) > 0 or s.get('charts', 0) > 0)

            eng_score = 100
            eng_issues = []

            if total_engagement < len(slides) * 0.3:
                eng_score -= 25
                eng_issues.append("Low engagement elements")

            if visual_slides < len(slides) * 0.4:
                eng_score -= 20
                eng_issues.append(f"Only {visual_slides}/{len(slides)} slides have visuals")

            question_slides = sum(1 for s in slides if '?' in ' '.join(s.get('text_content', [])))
            if question_slides < len(slides) * 0.2:
                eng_score -= 15
                eng_issues.append("Few slides have questions")

            scores['student_engagement'] = {
                'overall_score': max(0, eng_score),
                'issues': eng_issues,
                'recommendations': ['Add more interactive elements'] if eng_issues else [],
                'metrics': {
                    'engagement_elements': total_engagement,
                    'visual_slides': visual_slides,
                    'question_slides': question_slides,
                    'total_slides': len(slides)
                }
            }
        else:
            scores['student_engagement'] = {
                'overall_score': 50, 'issues': [], 'recommendations': [], 'metrics': {}
            }

        # 5. TECHNICAL DEPTH
        if slides:
            total_tech = sum(s.get('technical_terms', 0) for s in slides)
            avg_tech = total_tech / len(slides) if slides else 0

            tech_score = 100
            tech_issues = []

            if avg_tech < 1:
                tech_score -= 30
                tech_issues.append("Low technical depth for B.Tech level")
            elif avg_tech < 2:
                tech_score -= 15
                tech_issues.append("Could have more technical detail")

            # Check for formulas in technical subjects
            formula_count = presentation_data.get('educational_elements', {}).get('formulas', 0)
            subject = st.session_state.get('subject_area', '')
            if subject in ['Mathematics', 'Physics'] and formula_count == 0:
                tech_score -= 20
                tech_issues.append(f"No formulas found for {subject}")

            scores['technical_depth'] = {
                'overall_score': max(0, min(100, tech_score)),
                'issues': tech_issues,
                'recommendations': [],
                'metrics': {
                    'total_technical_terms': total_tech,
                    'avg_per_slide': round(avg_tech, 2),
                    'formulas': formula_count
                }
            }
        else:
            scores['technical_depth'] = {
                'overall_score': 50, 'issues': [], 'recommendations': [], 'metrics': {}
            }

        # 6. VISUAL DESIGN
        if slides:
            word_counts = [s.get('word_count', 0) for s in slides]
            avg_words = statistics.mean(word_counts) if word_counts else 0

            design_score = 100
            design_issues = []

            # Word count optimization
            if avg_words > 60:
                design_score -= 25
                design_issues.append(f"Slides too text-heavy (avg {avg_words:.0f} words)")
            elif avg_words > 80:
                design_score -= 40

            # Text-only slides
            text_only = sum(1 for s in slides
                            if s.get('images', 0) == 0 and s.get('charts', 0) == 0)
            if text_only > len(slides) * 0.6:
                design_score -= 20
                design_issues.append(f"{text_only}/{len(slides)} slides lack visual elements")

            # Overcrowded slides
            crowded = sum(1 for s in slides if s.get('word_count', 0) > 100)
            if crowded > 0:
                design_score -= crowded * 5
                design_issues.append(f"{crowded} slides are overcrowded")

            scores['visual_design'] = {
                'overall_score': max(0, design_score),
                'issues': design_issues,
                'recommendations': [],
                'metrics': {
                    'avg_word_count': round(avg_words, 1),
                    'text_only_slides': text_only,
                    'overcrowded_slides': crowded
                }
            }
        else:
            scores['visual_design'] = {
                'overall_score': 50, 'issues': [], 'recommendations': [], 'metrics': {}
            }

        # 7. LEARNING PROGRESSION
        if slides and len(slides) >= 3:
            # Check for introduction
            has_intro = any(
                any(kw in ' '.join(s.get('text_content', [])).lower()
                    for kw in ['introduction', 'overview', 'agenda', 'objectives'])
                for s in slides[:3]
            )

            # Check for conclusion
            has_conclusion = any(
                any(kw in ' '.join(s.get('text_content', [])).lower()
                    for kw in ['conclusion', 'summary', 'recap', 'takeaway'])
                for s in slides[-3:]
            )

            prog_score = 70
            prog_issues = []

            if has_intro:
                prog_score += 15
            else:
                prog_issues.append("Missing clear introduction/objectives")

            if has_conclusion:
                prog_score += 15
            else:
                prog_issues.append("Missing conclusion/summary")

            # Content flow (check word count variance)
            if len(word_counts) > 3:
                variance = statistics.variance(word_counts)
                if variance > 2000:  # High variance = inconsistent
                    prog_score -= 10
                    prog_issues.append("Inconsistent content distribution")

            scores['learning_progression'] = {
                'overall_score': min(100, prog_score),
                'issues': prog_issues,
                'recommendations': [],
                'metrics': {
                    'has_introduction': has_intro,
                    'has_conclusion': has_conclusion
                }
            }
        else:
            scores['learning_progression'] = {
                'overall_score': 50, 'issues': [], 'recommendations': [], 'metrics': {}
            }

        return scores

    def run_complete_analysis(self, file_path: Path,
                              use_gemini: bool = False, gemini_key: str = "", gemini_model: str = "",
                              use_groq: bool = False, groq_key: str = "", groq_model: str = "") -> Dict[str, Any]:
        """Run the COMPLETE analysis pipeline"""
        logger.info("Starting complete analysis")

        # Extract presentation data
        presentation_data = self.extract_presentation_data(file_path)

        if 'error' in presentation_data:
            raise Exception(presentation_data['error'])

        # Store total slides for reference
        st.session_state['total_slides'] = presentation_data['total_slides']

        # AI Analysis
        ai_results = {}

        if use_gemini and gemini_key and GEMINI_AVAILABLE:
            ai_results = self.run_gemini_analysis(
                file_path, presentation_data, gemini_key, gemini_model
            )
        elif use_groq and groq_key and LANGCHAIN_AVAILABLE:
            ai_results = self.run_groq_analysis(
                presentation_data, groq_key, groq_model
            )

        # Calculate quality scores
        quality_scores = self.calculate_quality_scores(presentation_data, ai_results)

        # Calculate weighted overall score
        weights = {
            'conceptual_accuracy': 0.25,
            'visual_alignment': 0.15,
            'educational_content': 0.20,
            'student_engagement': 0.15,
            'technical_depth': 0.10,
            'visual_design': 0.10,
            'learning_progression': 0.05
        }

        overall_score = sum(
            quality_scores[cat]['overall_score'] * weight
            for cat, weight in weights.items()
        )

        # Compile results
        results = {
            'overall_score': round(overall_score, 1),
            'enhanced_data': presentation_data,
            'quality_analysis': quality_scores,
            'ai_results': ai_results,
            'analysis_timestamp': datetime.now().isoformat(),
            'file_name': file_path.name,
            'target_audience': 'B.Tech Engineering Students',
            'subject_area': st.session_state.get('subject_area', 'Engineering'),
            'analysis_type': 'Complete Multimodal Analysis'
        }

        # FINAL CLEANUP of temp file
        try:
            os.remove(file_path)
            clean_memory()
        except:
            pass

        logger.info(f"Analysis complete. Score: {overall_score}")
        return results


# ============================================
# VISUALIZATION FUNCTIONS (THEME AWARE)
# ============================================

def create_quality_visualizations(results: Dict[str, Any]):
    """Create comprehensive visualizations - THEME AWARE"""
    st.markdown("### 📊 Quality Analysis Visualizations")

    slides = results.get('enhanced_data', {}).get('slides', [])
    quality = results.get('quality_analysis', {})

    if not slides:
        st.warning("No slide data for visualization")
        return

    # Use default matplotlib style which usually adapts better or specific dark mode friendly style
    plt.style.use('default')

    # Create figure with transparent background for theme compatibility
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    fig.patch.set_alpha(0.0)  # Transparent figure background

    # 1. Category Scores Bar Chart
    ax1 = axes[0, 0]
    ax1.patch.set_alpha(0.0)  # Transparent axis background
    categories = list(quality.keys())
    scores = [quality[cat]['overall_score'] for cat in categories]
    colors = ['#10B981' if s >= 90 else '#EF4444' for s in scores]

    y_pos = range(len(categories))
    bars = ax1.barh(y_pos, scores, color=colors)
    ax1.set_yticks(y_pos)
    ax1.set_yticklabels([cat.replace('_', ' ').title() for cat in categories])
    ax1.set_xlabel('Score')
    ax1.set_title('Quality Scores by Category')
    ax1.set_xlim(0, 100)

    # Set text colors to work in both themes (auto) or specific grey
    ax1.tick_params(axis='x', colors='gray')
    ax1.tick_params(axis='y', colors='gray')
    ax1.title.set_color('gray')
    ax1.xaxis.label.set_color('gray')

    for bar, score in zip(bars, scores):
        ax1.text(bar.get_width() + 1, bar.get_y() + bar.get_height() / 2,
                 f'{score}', va='center', fontweight='bold', color='gray')

    # 2. Educational Elements Pie Chart
    ax2 = axes[0, 1]
    ax2.patch.set_alpha(0.0)
    edu = results['enhanced_data'].get('educational_elements', {})
    labels = [k.replace('_', ' ').title() for k in edu.keys()]
    values = list(edu.values())

    if sum(values) > 0:
        colors_pie = plt.cm.Set3(np.linspace(0, 1, len(labels)))
        wedges, texts, autotexts = ax2.pie(values, labels=labels, autopct='%1.0f%%',
                                           colors=colors_pie, startangle=90)
        ax2.set_title('Educational Elements Distribution', color='gray')
        for text in texts: text.set_color('gray')
        for autotext in autotexts: autotext.set_color('black')  # Keep inside pie readable
    else:
        ax2.text(0.5, 0.5, 'No Educational\nElements Found',
                 ha='center', va='center', fontsize=12, color='gray')
        ax2.set_title('Educational Elements Distribution', color='gray')

    # 3. Slide Scores Line Chart
    ax3 = axes[1, 0]
    ax3.patch.set_alpha(0.0)
    if results.get('ai_results'):
        slide_nums = sorted(results['ai_results'].keys())
        slide_scores = [results['ai_results'][n].get('overall_score', 50) for n in slide_nums]
        point_colors = ['#10B981' if s >= 90 else '#EF4444' for s in slide_scores]

        ax3.plot(slide_nums, slide_scores, '-', color='#888888', linewidth=1, alpha=0.5)
        for x, y, c in zip(slide_nums, slide_scores, point_colors):
            ax3.plot(x, y, 'o', color=c, markersize=8)

        ax3.axhline(y=90, color='#10B981', linestyle='--', label='Excellent (90)')
        ax3.fill_between(slide_nums, slide_scores, alpha=0.1, color='#7C3AED')
    else:
        slide_nums = [s['slide_number'] for s in slides]
        slide_scores = [s.get('basic_score', 50) for s in slides]
        point_colors = ['#10B981' if s >= 90 else '#EF4444' for s in slide_scores]

        ax3.plot(slide_nums, slide_scores, '-', color='#888888', linewidth=1, alpha=0.5)
        for x, y, c in zip(slide_nums, slide_scores, point_colors):
            ax3.plot(x, y, 'o', color=c, markersize=8)

    ax3.set_xlabel('Slide Number', color='gray')
    ax3.set_ylabel('Quality Score', color='gray')
    ax3.set_title('Slide-by-Slide Quality Scores', color='gray')
    ax3.tick_params(axis='x', colors='gray')
    ax3.tick_params(axis='y', colors='gray')
    ax3.set_ylim(0, 100)
    ax3.grid(True, alpha=0.3)

    # 4. Word Count Distribution
    ax4 = axes[1, 1]
    ax4.patch.set_alpha(0.0)
    word_counts = [s['word_count'] for s in slides]
    slide_nums = [s['slide_number'] for s in slides]

    colors_bars = ['#10B981' if 20 <= wc <= 50 else '#EF4444' for wc in word_counts]
    ax4.bar(slide_nums, word_counts, color=colors_bars)
    ax4.axhline(y=50, color='#10B981', linestyle='--', label='Optimal Max (50)')
    ax4.axhline(y=80, color='#EF4444', linestyle='--', label='Too Much (80)')
    ax4.set_xlabel('Slide Number', color='gray')
    ax4.set_ylabel('Word Count', color='gray')
    ax4.set_title('Word Count per Slide', color='gray')
    ax4.tick_params(axis='x', colors='gray')
    ax4.tick_params(axis='y', colors='gray')
    ax4.legend(facecolor='white', framealpha=0.5)  # Semi-transparent legend

    plt.tight_layout()
    st.pyplot(fig)

    # Close plots to free memory
    plt.close(fig)

    # Word Cloud
    st.markdown("### ☁️ Content Word Cloud")
    all_text = ' '.join([' '.join(s.get('text_content', [])) for s in slides])

    if len(all_text) > 100:
        try:
            stop_words = set(stopwords.words('english'))
            words = [w for w in all_text.lower().split()
                     if w not in stop_words and len(w) > 3 and w.isalpha()]

            if words:
                # Wordcloud handles its own background
                wordcloud = WordCloud(
                    width=1000, height=400,
                    background_color='black',  # Dark background usually looks better
                    colormap='viridis',
                    max_words=100
                ).generate(' '.join(words))

                fig_wc, ax_wc = plt.subplots(figsize=(12, 5))
                fig_wc.patch.set_alpha(0.0)
                ax_wc.imshow(wordcloud, interpolation='bilinear')
                ax_wc.axis('off')
                ax_wc.set_title('Key Terms and Concepts', fontsize=14, fontweight='bold', color='gray')
                st.pyplot(fig_wc)
                plt.close(fig_wc)
        except Exception as e:
            st.warning(f"Could not generate word cloud: {e}")


# ============================================
# DISPLAY FUNCTIONS (THEME AWARE)
# ============================================

def display_overall_dashboard(results: Dict[str, Any]):
    """Display main quality dashboard - THEME AWARE"""
    st.markdown("## 📊 Quality Assessment Dashboard")

    score = results['overall_score']

    # Use generic border colors, avoid hardcoded white/black backgrounds
    if score >= 90:
        rating = "🏆 EXCELLENT"
        border = "#10B981"  # Green
    else:
        rating = "⚠️ NEEDS IMPROVEMENT"
        border = "#EF4444"  # Red

    # Use Streamlit's native CSS variables for background/text adaptation
    st.markdown(f"""
    <div style="
        border: 3px solid {border}; 
        background-color: var(--secondary-background-color);
        border-radius: 20px; 
        padding: 2rem; 
        text-align: center; 
        margin: 1rem 0;
        color: var(--text-color);
    ">
        <h2 style="margin: 0; color: var(--text-color);">Overall Educational Quality Score</h2>
        <h1 style="font-size: 4rem; margin: 0.5rem 0; color: {border};">{score}/100</h1>
        <h3 style="margin: 0; color: var(--text-color);">{rating}</h3>
        <p style="margin-top: 1rem; color: var(--text-color);">
            <strong>Target:</strong> {results.get('target_audience', 'B.Tech Students')} | 
            <strong>Subject:</strong> {results.get('subject_area', 'Engineering')}
        </p>
    </div>
    """, unsafe_allow_html=True)

    # Quick stats using native metrics (theme aware by default)
    data = results.get('enhanced_data', {})
    col1, col2, col3, col4, col5 = st.columns(5)

    with col1:
        st.metric("📖 Total Slides", data.get('total_slides', 0))
    with col2:
        st.metric("💡 Examples", data.get('educational_elements', {}).get('examples', 0))
    with col3:
        st.metric("❓ Questions", data.get('educational_elements', {}).get('questions', 0))
    with col4:
        st.metric("🖼️ Images", data.get('images_count', 0))
    with col5:
        st.metric("📊 Charts", data.get('charts_count', 0))

    # Category breakdown
    st.markdown("### 📈 Quality by Category")

    quality = results.get('quality_analysis', {})
    cols = st.columns(4)

    for i, (cat, analysis) in enumerate(quality.items()):
        col = cols[i % 4]
        cat_score = analysis.get('overall_score', 0)

        if cat_score >= 90:
            emoji = "🟢"
            color = "#10B981"
        else:
            emoji = "🔴"
            color = "#EF4444"

        with col:
            cat_name = UltimateEduPPTAnalyzer.QUALITY_CATEGORIES.get(cat, cat.replace('_', ' ').title())
            # Theme aware card
            st.markdown(f"""
            <div style="
                background-color: var(--secondary-background-color); 
                padding: 1rem; 
                border-radius: 10px; 
                border-left: 4px solid {color}; 
                margin: 0.5rem 0; 
                box-shadow: 0 2px 5px rgba(0,0,0,0.1);
                color: var(--text-color);
            ">
                <div style="font-size: 0.9rem; opacity: 0.8;">{emoji} {cat_name}</div>
                <div style="font-size: 1.5rem; font-weight: bold; color: {color};">{cat_score}/100</div>
            </div>
            """, unsafe_allow_html=True)


def display_slide_analysis(results: Dict[str, Any]):
    """Display detailed slide-by-slide analysis"""
    ai_results = results.get('ai_results', {})

    if not ai_results:
        st.info("📝 Enable AI analysis (Gemini or Groq) for detailed slide-by-slide recommendations.")

        # Show basic analysis from presentation data
        slides = results.get('enhanced_data', {}).get('slides', [])
        if slides:
            st.markdown("### 📄 Basic Slide Overview")
            for slide in slides:
                with st.expander(f"Slide {slide['slide_number']} - Basic Score: {slide.get('basic_score', 50)}/100"):
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        st.metric("Words", slide['word_count'])
                        st.metric("Images", slide['images'])
                    with col2:
                        st.metric("Density", slide['content_density'].title())
                        st.metric("Purpose", slide['slide_purpose'].title())
                    with col3:
                        st.metric("Technical", slide['technical_terms'])
                        st.metric("Engagement", slide['engagement_elements'])
        return

    st.markdown("## 🔬 Detailed Slide-by-Slide Analysis")

    # Summary
    total = len(ai_results)
    critical = sum(1 for r in ai_results.values() if r.get('priority_level') == 'CRITICAL')
    high = sum(1 for r in ai_results.values() if r.get('priority_level') == 'HIGH')

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Total Analyzed", total)
    with col2:
        st.metric("Critical Issues", critical, delta="Urgent!" if critical else None)
    with col3:
        st.metric("High Priority", high)
    with col4:
        avg_score = statistics.mean([r.get('overall_score', 50) for r in ai_results.values()])
        st.metric("Avg Score", f"{avg_score:.0f}")

    # Filter
    st.markdown("### 🔍 Filter by Priority")
    filter_opt = st.selectbox(
        "Show:",
        ["All Slides", "CRITICAL Only", "HIGH Only", "MEDIUM Only", "LOW Only"],
        key="slide_filter_select"
    )

    if filter_opt != "All Slides":
        priority = filter_opt.replace(" Only", "")
        filtered = {k: v for k, v in ai_results.items() if v.get('priority_level') == priority}
    else:
        filtered = ai_results

    if not filtered:
        st.info(f"No slides with {filter_opt.lower()}")
        return

    # Display slides
    for slide_num in sorted(filtered.keys()):
        r = filtered[slide_num]
        priority = r.get('priority_level', 'MEDIUM')
        score = r.get('overall_score', 50)

        # Priority colors (CSS safe)
        if priority == 'CRITICAL':
            border_color = '#EF4444'  # Red
            badge_bg = 'rgba(239, 68, 68, 0.2)'
        elif priority == 'HIGH':
            border_color = '#F97316'  # Orange
            badge_bg = 'rgba(249, 115, 22, 0.2)'
        elif priority == 'LOW':
            border_color = '#22C55E'  # Green
            badge_bg = 'rgba(34, 197, 94, 0.2)'
        else:  # MEDIUM
            border_color = '#EAB308'  # Yellow
            badge_bg = 'rgba(234, 179, 8, 0.2)'

        with st.expander(f"📄 Slide {slide_num} | Score: {score}/100 | Priority: {priority}",
                         expanded=(priority in ['CRITICAL', 'HIGH'])):

            # Priority badge
            st.markdown(f"""
            <div style="background-color: {badge_bg}; border: 1px solid {border_color}; color: {border_color}; 
                        padding: 0.3rem 1rem; border-radius: 20px; font-weight: bold; font-size: 0.85rem; 
                        display: inline-block; margin-bottom: 10px;">
                {priority} PRIORITY
            </div>
            """, unsafe_allow_html=True)

            # Scores
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Overall", f"{score}/100")
            with col2:
                st.metric("Conceptual", f"{r.get('conceptual_accuracy_score', '-')}/100")
            with col3:
                st.metric("Visual-Text", f"{r.get('visual_text_alignment_score', '-')}/100")
            with col4:
                st.metric("Educational", f"{r.get('educational_effectiveness_score', '-')}/100")

            # Critical Errors
            if r.get('critical_errors') and any(r['critical_errors']):
                st.error("🚨 **Critical Errors:**")
                for err in r['critical_errors']:
                    if err and len(err) > 3:
                        st.markdown(f"- {err}")

            # Immediate Fixes
            if r.get('immediate_fixes') and any(r['immediate_fixes']):
                st.warning("⚡ **Immediate Fixes Required:**")
                for i, fix in enumerate(r['immediate_fixes'], 1):
                    if fix and len(fix) > 3:
                        st.markdown(f"{i}. {fix}")

            # Other issues
            tabs = st.tabs(["Conceptual", "Visual", "Text", "Recommendations", "Positives"])

            with tabs[0]:
                if r.get('conceptual_issues') and len(r['conceptual_issues']) > 5:
                    st.warning(r['conceptual_issues'])
                else:
                    st.success("No major conceptual issues")

            with tabs[1]:
                if r.get('visual_issues') and len(r['visual_issues']) > 5:
                    st.info(r['visual_issues'])
                if r.get('alignment_problems') and len(r['alignment_problems']) > 5:
                    st.warning(f"**Alignment Issues:** {r['alignment_problems']}")

            with tabs[2]:
                if r.get('text_issues') and len(r['text_issues']) > 5:
                    st.info(r['text_issues'])
                else:
                    st.success("No text issues found")

            with tabs[3]:
                if r.get('detailed_recommendations'):
                    st.markdown(r['detailed_recommendations'])
                if r.get('student_impact'):
                    st.markdown(f"**Student Impact:** {r['student_impact']}")

            with tabs[4]:
                if r.get('positive_aspects') and len(r['positive_aspects']) > 5:
                    st.success(r['positive_aspects'])
                else:
                    st.info("Review other tabs for analysis")


def display_issues_summary(results: Dict[str, Any]):
    """Display summary of all issues - THEME AWARE"""
    st.markdown("## 🚨 Issues Summary")

    quality = results.get('quality_analysis', {})

    all_issues = []
    for cat, analysis in quality.items():
        cat_name = UltimateEduPPTAnalyzer.QUALITY_CATEGORIES.get(cat, cat)
        for issue in analysis.get('issues', []):
            if issue:
                all_issues.append({'category': cat_name, 'issue': issue})

    if not all_issues:
        st.success("🎉 No major issues found! The presentation meets quality standards.")
        return

    st.markdown(f"**Found {len(all_issues)} issues to address:**")

    for i, item in enumerate(all_issues, 1):
        # Use st.warning for theme-safe styling instead of raw HTML with hardcoded colors
        st.warning(f"**#{i} [{item['category']}]**\n\n{item['issue']}")


def generate_report(results: Dict[str, Any]) -> str:
    """Generate comprehensive markdown report"""
    timestamp = datetime.fromisoformat(results['analysis_timestamp']).strftime('%Y-%m-%d %H:%M:%S')

    report = f"""# Educational PowerPoint Quality Analysis Report

## Summary
- **File:** {results['file_name']}
- **Target Audience:** {results['target_audience']}
- **Subject Area:** {results.get('subject_area', 'Engineering')}
- **Analysis Date:** {timestamp}
- **Overall Score:** {results['overall_score']}/100

---

## Quality Scores

| Category | Score | Status |
|----------|-------|--------|
"""

    for cat, analysis in results.get('quality_analysis', {}).items():
        score = analysis['overall_score']
        # Updated report status logic
        status = '✅' if score >= 90 else '❌'
        cat_name = UltimateEduPPTAnalyzer.QUALITY_CATEGORIES.get(cat, cat)
        report += f"| {cat_name} | {score}/100 | {status} |\n"

    # Educational elements
    edu = results.get('enhanced_data', {}).get('educational_elements', {})
    report += f"""
---

## Educational Elements Found

| Element | Count |
|---------|-------|
| Examples | {edu.get('examples', 0)} |
| Definitions | {edu.get('definitions', 0)} |
| Questions | {edu.get('questions', 0)} |
| Code Blocks | {edu.get('code_blocks', 0)} |
| Diagrams | {edu.get('diagrams', 0)} |
| Formulas | {edu.get('formulas', 0)} |

---

## Issues Found
"""

    for cat, analysis in results.get('quality_analysis', {}).items():
        if analysis.get('issues'):
            cat_name = UltimateEduPPTAnalyzer.QUALITY_CATEGORIES.get(cat, cat)
            report += f"\n### {cat_name}\n"
            for issue in analysis['issues']:
                report += f"- ❌ {issue}\n"

    # Slide analysis
    if results.get('ai_results'):
        report += "\n---\n\n## Slide-by-Slide Analysis\n"

        for slide_num in sorted(results['ai_results'].keys()):
            r = results['ai_results'][slide_num]
            report += f"""
### Slide {slide_num}
- **Score:** {r.get('overall_score', 'N/A')}/100
- **Priority:** {r.get('priority_level', 'N/A')}
"""
            if r.get('critical_errors'):
                report += "\n**Critical Errors:**\n"
                for err in r['critical_errors']:
                    if err:
                        report += f"- {err}\n"

            if r.get('immediate_fixes'):
                report += "\n**Immediate Fixes:**\n"
                for fix in r['immediate_fixes']:
                    if fix:
                        report += f"- {fix}\n"

    report += f"""
---

*Report generated by EduPPT Quality Analyzer Ultimate v7.7*
*Analysis powered by Google Gemini / Groq Llama*
"""

    return report


# ============================================
# STREAMLIT UI STYLING (DARK/LIGHT FIX)
# ============================================

def apply_custom_css():
    """Apply custom CSS styling using CSS Variables for Theme Adaptation"""
    st.markdown("""
    <style>
    /* Main header - Gradient stays, but text is forced white for readability */
    .main-header {
        background: linear-gradient(135deg, #1E40AF 0%, #7C3AED 50%, #DB2777 100%);
        padding: 2.5rem;
        border-radius: 20px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
        box-shadow: 0 15px 50px rgba(30, 64, 175, 0.4);
    }

    /* Cards - Adapt to theme background */
    .quality-card {
        background-color: var(--secondary-background-color);
        padding: 1.5rem;
        border-radius: 15px;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.1);
        margin: 1rem 0;
        border-left: 5px solid #7C3AED;
        color: var(--text-color);
    }

    /* Feature badges */
    .feature-badge {
        background: linear-gradient(135deg, #7C3AED 0%, #DB2777 100%);
        color: white;
        padding: 0.3rem 0.8rem;
        border-radius: 15px;
        font-size: 0.85rem;
        font-weight: bold;
        margin: 0.2rem;
        display: inline-block;
    }

    .gemini-badge {
        background: linear-gradient(135deg, #4285F4 0%, #34A853 33%, #FBBC05 66%, #EA4335 100%);
        color: white;
        padding: 0.5rem 1.5rem;
        border-radius: 25px;
        font-weight: bold;
        display: inline-block;
        margin: 0.5rem;
    }

    /* Buttons */
    .stButton > button {
        background: linear-gradient(135deg, #7C3AED 0%, #DB2777 100%);
        color: white;
        border: none;
        border-radius: 12px;
        padding: 0.75rem 2rem;
        font-weight: bold;
        font-size: 1rem;
        transition: all 0.3s ease;
    }

    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 25px rgba(124, 58, 237, 0.4);
    }
    </style>
    """, unsafe_allow_html=True)


# ============================================
# CLEANUP UTILS (STARTUP FIX)
# ============================================
def cleanup_data_dir():
    """Wipe the data directory on startup to remove old crash files"""
    data_dir = Path("data")
    if data_dir.exists():
        for file in data_dir.glob("*"):
            try:
                if file.is_file():
                    file.unlink()
                    logger.info(f"🧹 Startup cleanup: Removed {file.name}")
            except Exception as e:
                logger.warning(f"Could not clean {file}: {e}")


# ============================================
# MAIN APPLICATION
# ============================================

def main():
    """Main application entry point"""

    # --- STARTUP CLEANUP ---
    if 'startup_cleanup_done' not in st.session_state:
        cleanup_data_dir()
        st.session_state['startup_cleanup_done'] = True
    # -----------------------

    # Check critical dependencies
    dep_check = st.session_state.dep_check
    if dep_check['missing']:
        st.error("❌ Missing Critical Dependencies")
        for err in dep_check['missing']:
            st.error(err)
        st.code('\n'.join(dep_check['install_commands']), language='bash')
        st.stop()

    # Apply styling
    apply_custom_css()

    # Header
    st.markdown("""
    <div class="main-header">
        <h1>🎓 EduPPT Quality Analyzer Ultimate</h1>
        <p>Complete Multimodal Analysis for Educational Presentations</p>
        <div>
            <span class="gemini-badge">✨ Google Gemini 2.0</span>
            <span class="feature-badge">🦙 Groq Llama 4</span>
            <span class="feature-badge">📊 7 Quality Categories</span>
            <span class="feature-badge">🔬 Slide-by-Slide AI</span>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Warnings for optional dependencies
    if dep_check['warnings']:
        with st.expander("⚠️ Optional Dependencies"):
            for warn in dep_check['warnings']:
                st.warning(warn)

    # Initialize analyzer
    if 'analyzer' not in st.session_state:
        with st.spinner("Initializing analyzer..."):
            st.session_state.analyzer = UltimateEduPPTAnalyzer()

    analyzer = st.session_state.analyzer

    # Sidebar configuration
    with st.sidebar:
        st.header("⚙️ Analysis Configuration")

        st.subheader("🎯 Target Settings")
        st.session_state['subject_area'] = st.selectbox(
            "Subject Area",
            ["Computer Science", "Electronics", "Mechanical", "Civil",
             "Mathematics", "Physics", "Chemistry", "General Engineering",
             "Data Science", "Artificial Intelligence", "Networking"]
        )

        st.markdown("---")

        st.subheader("🤖 AI Analysis Engine")

        ai_choice = st.radio(
            "Select AI Model",
            ["Google Gemini (Recommended)", "Groq Llama", "No AI (Basic Only)"],
            help="Gemini provides best multimodal analysis"
        )

        use_gemini = ai_choice == "Google Gemini (Recommended)"
        use_groq = ai_choice == "Groq Llama"

        gemini_key = ""
        gemini_model = "gemini-2.0-flash"
        groq_key = ""
        groq_model = "llama-3.3-70b-versatile"

        if use_gemini:
            if not GEMINI_AVAILABLE:
                st.error("google-generativeai not installed")
                st.code("pip install google-generativeai")
            else:
                # Check ENV for key
                env_gemini_key = os.getenv("GOOGLE_API_KEY")

                if env_gemini_key:
                    st.success("✅ Gemini Key loaded from .env")
                    gemini_key = env_gemini_key
                else:
                    gemini_key = st.text_input(
                        "Google API Key",
                        type="password",
                        help="Get from https://makersuite.google.com/app/apikey"
                    )

                gemini_model = st.selectbox(
                    "Gemini Model",
                    options=list(GeminiMultimodalAnalyzer.MODELS.keys()),
                    format_func=lambda x: GeminiMultimodalAnalyzer.MODELS[x]
                )

                if gemini_key:
                    st.success("✅ Gemini Ready")
                else:
                    st.warning("⚠️ Enter API key or set GOOGLE_API_KEY in .env")

        elif use_groq:
            if not LANGCHAIN_AVAILABLE:
                st.error("langchain-groq not installed")
                st.code("pip install langchain-groq langchain-core")
            else:
                # Check ENV for key
                env_groq_key = os.getenv("GROQ_API_KEY")

                if env_groq_key:
                    st.success("✅ Groq Key loaded from .env")
                    groq_key = env_groq_key
                else:
                    groq_key = st.text_input(
                        "Groq API Key",
                        type="password",
                        help="Get from https://console.groq.com"
                    )

                groq_model = st.selectbox(
                    "Groq Model",
                    options=list(GroqLlamaAnalyzer.MODELS.keys()),
                    format_func=lambda x: GroqLlamaAnalyzer.MODELS[x]
                )

                if groq_key:
                    st.success("✅ Groq Ready")
                else:
                    st.warning("⚠️ Enter API key or set GROQ_API_KEY in .env")

        st.markdown("---")
        st.markdown("""
        ### 📋 Analysis Features
        - ✅ Conceptual accuracy check
        - ✅ Visual-text alignment
        - ✅ Educational content quality
        - ✅ Student engagement scoring
        - ✅ Technical depth analysis
        - ✅ Visual design evaluation
        - ✅ Learning progression check
        - ✅ Detailed visualizations
        - ✅ Comprehensive reports
        """)

    # Main content area
    col1, col2 = st.columns([2, 1])

    with col1:
        st.markdown("""
        <div class="quality-card">
            <h3>📁 Upload Your Presentation</h3>
            <p>Upload a PowerPoint file (.pptx) for comprehensive quality analysis</p>
            <p><small>Works for ANY educational content - universal analyzer</small></p>
        </div>
        """, unsafe_allow_html=True)

        uploaded_file = st.file_uploader(
            "Choose a PowerPoint file",
            type=["pptx"],
            help="Upload your .pptx file for analysis"
        )

        if uploaded_file:
            file_size = uploaded_file.size / (1024 * 1024)
            st.info(f"📄 **{uploaded_file.name}** ({file_size:.2f} MB)")

            if st.button("🚀 Start Complete Analysis", type="primary"):
                try:
                    # Clear session state for memory
                    clean_memory()

                    # Save file
                    with st.spinner("Saving file..."):
                        save_path = analyzer.save_uploaded_file(uploaded_file)

                    # Run analysis
                    results = analyzer.run_complete_analysis(
                        save_path,
                        use_gemini=use_gemini and bool(gemini_key),
                        gemini_key=gemini_key,
                        gemini_model=gemini_model,
                        use_groq=use_groq and bool(groq_key),
                        groq_key=groq_key,
                        groq_model=groq_model
                    )

                    st.session_state.results = results
                    st.success(f"✅ Analysis complete! Score: {results['overall_score']}/100")

                except Exception as e:
                    st.error(f"❌ Error: {str(e)}")
                    logger.error(f"Analysis error: {traceback.format_exc()}")

    with col2:
        st.markdown("""
        <div class="quality-card">
            <h3>🌟 Why This Analyzer?</h3>
            <ul>
                <li><strong>Universal:</strong> Works for ANY content</li>
                <li><strong>Accurate:</strong> Gemini 2.0 multimodal AI</li>
                <li><strong>Complete:</strong> 7 quality dimensions</li>
                <li><strong>Detailed:</strong> Slide-by-slide analysis</li>
                <li><strong>Actionable:</strong> Specific fix recommendations</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)

    # Display results
    if 'results' in st.session_state:
        results = st.session_state.results

        st.markdown("---")

        # Main dashboard
        display_overall_dashboard(results)

        # Visualizations
        create_quality_visualizations(results)

        # Slide analysis
        display_slide_analysis(results)

        # Issues summary
        display_issues_summary(results)

        # Download section
        st.markdown("---")
        st.markdown("### 📥 Download Reports")

        col1, col2 = st.columns(2)

        with col1:
            report = generate_report(results)
            st.download_button(
                "📄 Download Report (Markdown)",
                data=report,
                file_name=f"quality_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
                mime="text/markdown",
                type="primary"
            )

        with col2:
            # JSON summary
            json_data = {
                'overall_score': results['overall_score'],
                'file_name': results['file_name'],
                'subject_area': results.get('subject_area'),
                'total_slides': results['enhanced_data']['total_slides'],
                'category_scores': {
                    cat: analysis['overall_score']
                    for cat, analysis in results['quality_analysis'].items()
                },
                'educational_elements': results['enhanced_data']['educational_elements'],
                'slides_analyzed': len(results.get('ai_results', {}))
            }

            st.download_button(
                "📊 Download Data (JSON)",
                data=json.dumps(json_data, indent=2),
                file_name=f"analysis_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                mime="application/json"
            )

    # Footer
    st.markdown("---")
    st.markdown("""
    <div style="text-align: center; color: gray; padding: 2rem; 
                background: linear-gradient(135deg, var(--secondary-background-color) 0%, var(--background-color) 100%); 
                border-radius: 15px;">
        <h4>🎓 EduPPT Quality Analyzer Ultimate v7.7.0</h4>
        <p><strong>Complete Enterprise Edition</strong> - AWS EC2 Optimized</p>
        <p>Powered by Google Gemini 2.0 Flash & Groq Llama 4</p>
        <p>Specialized for B.Tech Engineering Education</p>
    </div>
    """, unsafe_allow_html=True)


if __name__ == "__main__":
    main()